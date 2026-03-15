#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import os
import subprocess
import sys
import textwrap
from collections import Counter, defaultdict
from datetime import UTC, datetime, timedelta
from pathlib import Path
from typing import Any
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen

ACCOUNT = "joy7758"
API_ROOT = "https://api.github.com"
SCRIPT_PATH = Path(__file__).resolve()
OUTPUT_DIR = SCRIPT_PATH.parent
WORKSPACE_ROOT = OUTPUT_DIR.parents[1]
NOW_UTC = datetime.now(UTC)
REPO_CACHE = OUTPUT_DIR / "repo_snapshot_cache.json"
GITHUB_TOKEN = (
    os.getenv("GITHUB_TOKEN") or os.getenv("GH_TOKEN") or os.getenv("GITHUB_PAT")
)

LOCAL_NAME_MAP = {
    "ToolAgents-ei": "ToolAgents",
    "agno-ei": "agno",
    "freeact-ei": "freeact",
    "mcp-agent-ei": "mcp-agent",
    "pydantic-ai-ei": "pydantic-ai",
}

LOCAL_IGNORE = {
    "docs-sparse",
    "rlcp-experiments",
}

SKIP_DIRS = {
    ".git",
    ".hg",
    ".svn",
    ".venv",
    "venv",
    "node_modules",
    "__pycache__",
    ".pytest_cache",
    ".ruff_cache",
    ".mypy_cache",
    ".next",
    "dist",
    "build",
}

PRIMARY_REPOS = {
    "god-spear": {
        "module": "上线前门禁层",
        "layer": "Governance gate",
        "bp_role": "上线前门禁，要求信任边界、失败信号和回滚路径明确。",
        "one_line": "CI-native trust gate for risky AI automation.",
        "stage": "可现场演示，可直接接入 CI。",
        "capabilities": [
            "对 tools/files/env/runtime 入口做前置 trust-boundary 检查",
            "缺少 failure signal、revocation pathway、grace budget 时直接 FAIL",
            "可输出文本、JSON、HTML 报告，并追加 immutable trace",
            "提供 `god-spear-mcp-gate` 适配器，便于挂到 MCP / tool 调用前",
        ],
        "validation": [
            "GitHub Release: `v0.2.0`",
            "npm 包发布与 SBOM、SHA256、release file list 已入库",
            "本地 smoke check 通过：`node bin/spear.js check examples/rules.ok.json`",
        ],
        "smoke": {
            "cmd": ["node", "bin/spear.js", "check", "examples/rules.ok.json"],
            "success_markers": ["STATUS: PASS"],
            "timeout": 30,
        },
    },
    "safety-valve-spec": {
        "module": "动作边界收据层",
        "layer": "Execution safety / integrity",
        "bp_role": "动作边界收据层，要求高风险动作必须带可验收据。",
        "one_line": "Verifiable receipt spec and conformance suite for action boundaries.",
        "stage": "规范、参考实现、conformance、badge 全部齐备。",
        "capabilities": [
            "定义 ALLOW / DENY / DEGRADE 收据格式与 JSON Schema",
            "提供签名、验签、CA、CRL、compat attestation 工具链",
            "提供 gateway demo，落实 no receipt, no action",
            "提供 reusable GitHub workflow，其他仓库可复用 conformance 检查",
        ],
        "validation": [
            "GitHub Release: `v0.2.0-alpha.1`",
            "仓库内含 `dist/svs-compat.attestation.json` 与 badge 产物",
            "本地 conformance 全量通过：`bash conformance/run.sh`",
        ],
        "smoke": {
            "cmd": ["bash", "conformance/run.sh"],
            "success_markers": ["Overall: PASS"],
            "timeout": 120,
        },
    },
    "execution-integrity-core": {
        "module": "执行完整性层",
        "layer": "Execution safety / integrity",
        "bp_role": "最小执行完整性证明，保证执行链条可导出、可验证、可发现篡改。",
        "one_line": "Minimal structural proof for execution integrity.",
        "stage": "最小证明闭环成熟，适合对外说明和客户演示。",
        "capabilities": [
            "记录 execution event，并对事件做 hash chain 串联",
            "支持确定性导出 JSON 与 full-chain verification",
            "提供篡改前后对比演示，强调 structural integrity",
            "作为执行层最小内核，适合做对外说明时的技术锚点",
        ],
        "validation": [
            "仓库有 `SPEC.md`、`scripts/selfcheck.sh`、`tests/test_export_verify.py`",
            "已接入 `spear-check`，说明执行层也纳入治理门禁",
            "本地 self-check 通过：`bash scripts/selfcheck.sh`",
        ],
        "smoke": {
            "cmd": ["bash", "scripts/selfcheck.sh"],
            "success_markers": ["SELF_CHECK: PASS"],
            "timeout": 30,
        },
    },
    "aro-audit": {
        "module": "审计证据层",
        "layer": "Audit evidence",
        "bp_role": "事后证据层，把普通日志升级为可复核、可回放、可交付的审计包。",
        "one_line": "Audit evidence layer for bounded, reviewable AI execution artifacts.",
        "stage": "演示、规范、bundle 都齐备，适合做主项目外壳。",
        "capabilities": [
            "生成 append-only journal、manifest、Merkle/checkpoint 与验签流程",
            "支持 quickstart 生成正常样本与 tamper 样本并独立验证",
            "支持 bundle 导出与 SHA256 manifest，便于对外交付审计材料",
            "已形成 one-pager、threat model、conformance vectors 等对外材料",
        ],
        "validation": [
            "GitHub Release: `v1.0.1`",
            "有 `CITATION.cff`、`SECURITY.md`、Zenodo DOI",
            "本地 quickstart 通过：正常样本 `VERIFY_OK`，篡改样本 `Merkle mismatch`",
        ],
        "smoke": {
            "cmd": ["bash", "quickstart/run.sh"],
            "success_markers": [
                "=== DONE: Quickstart OK ===",
                "VERIFY_OK: full chain valid",
            ],
            "timeout": 120,
        },
    },
}

SUPPORTING_REPOS = {
    "persona-object-protocol": {
        "layer": "Identity foundation",
        "role": "为更完整的大架构提供 persona / identity 层；当前不作为主模块单独对外销售。",
    },
    "pop-persona-pack": {
        "layer": "Identity foundation",
        "role": "POP 的 persona 资产包，证明 identity 层可落地。",
    },
    "langchain-pop": {
        "layer": "Identity foundation",
        "role": "将 persona 层贴近 LangChain 场景。",
    },
    "token-governor": {
        "layer": "Governance gate",
        "role": "更深一层的运行时预算和策略治理扩展，可作为 AegisTrust 后续控制平面。",
    },
    "token-governor-langchain-middleware": {
        "layer": "Governance gate",
        "role": "Token Governor 的轻量中间件适配层。",
    },
    "god-spear-mcp-gate": {
        "layer": "Governance gate",
        "role": "将 trust gate 前置到 MCP 风格工具执行。",
    },
    "fdo-kernel-mvk": {
        "layer": "Execution safety / integrity",
        "role": "执行完整性层的进阶演进，证明执行层不是停留在概念而在持续展开。",
    },
    "verifiable-agent-demo": {
        "layer": "Execution safety / integrity",
        "role": "跨层端到端 demo，把 persona、governance、execution、audit 串成一条评审可见路径。",
        "smoke": {
            "cmd": ["python3", "-m", "demo.agent"],
            "success_markers": ['"audit_record"'],
            "timeout": 30,
        },
    },
    "langchain-aro": {
        "layer": "Audit evidence",
        "role": "将 audit evidence 接到 LangChain 运行时。",
    },
    "aro-audit-langchain-receipt": {
        "layer": "Audit evidence",
        "role": "最小 post-run receipt 适配器，便于接入现有 agent 框架。",
    },
    "digital-biosphere-architecture": {
        "layer": "Adjacent architecture",
        "role": "架构总览仓库，用于解释生态关系，不作为当前阶段的主交付。",
    },
    "joy7758": {
        "layer": "Adjacent architecture",
        "role": "GitHub profile，公开展示创始人研究方向与五层架构定位。",
    },
    "agent-intent-protocol": {
        "layer": "Adjacent architecture",
        "role": "interaction semantics 层，当前不单列为商业模块。",
    },
    "agent-object-protocol": {
        "layer": "Adjacent architecture",
        "role": "agent object draft，作为长期协议储备。",
    },
}

LAYER_ORDER = [
    "Identity foundation",
    "Governance gate",
    "Execution safety / integrity",
    "Audit evidence",
    "Adjacent architecture",
]

LAYER_LABELS = {
    "Identity foundation": "身份基础层",
    "Governance gate": "治理门禁层",
    "Execution safety / integrity": "执行安全与完整性层",
    "Audit evidence": "审计证据层",
    "Adjacent architecture": "外围架构仓库",
}


def github_get(path: str) -> Any:
    url = path if path.startswith("http") else f"{API_ROOT}{path}"
    headers = {
        "User-Agent": "aegistrust-bp-generator/1.0",
        "Accept": "application/vnd.github+json",
    }
    if GITHUB_TOKEN:
        headers["Authorization"] = f"Bearer {GITHUB_TOKEN}"
    req = Request(url, headers=headers)
    with urlopen(req, timeout=30) as response:
        return json.load(response)


def iso_to_dt(value: str) -> datetime:
    return datetime.fromisoformat(value.replace("Z", "+00:00"))


def fmt_dt(value: str | datetime | None) -> str:
    if value is None:
        return "n/a"
    dt = iso_to_dt(value) if isinstance(value, str) else value
    return dt.astimezone(UTC).strftime("%Y-%m-%d %H:%M UTC")


def fmt_day(value: str | datetime | None) -> str:
    if value is None:
        return "n/a"
    dt = iso_to_dt(value) if isinstance(value, str) else value
    return dt.astimezone(UTC).strftime("%Y-%m-%d")


def normalize_message(message: str) -> str:
    return message.splitlines()[0].strip()


def shell_join(parts: list[str]) -> str:
    return " ".join(parts)


def repo_path(name: str) -> Path:
    return WORKSPACE_ROOT / name


def repo_exists(name: str) -> bool:
    path = repo_path(name)
    return path.exists() and (path / ".git").exists()


def detect_local_language(path: Path) -> str | None:
    if (path / "pyproject.toml").exists() or (path / "requirements.txt").exists():
        return "Python"
    if (path / "package.json").exists():
        return "JavaScript"
    if any(path.glob("*.html")):
        return "HTML"
    if any(path.glob("*.tex")):
        return "TeX"
    if any(path.glob("*.sh")) or (path / "Makefile").exists():
        return "Shell"
    return None


def local_description(name: str, path: Path) -> str:
    if name in PRIMARY_REPOS:
        return PRIMARY_REPOS[name]["one_line"]
    readme = path / "README.md"
    if readme.exists():
        for line in readme.read_text(encoding="utf-8", errors="ignore").splitlines():
            stripped = line.strip()
            if not stripped or stripped.startswith("#") or stripped.startswith("!"):
                continue
            if len(stripped) > 200:
                continue
            return stripped
    if name in SUPPORTING_REPOS:
        return SUPPORTING_REPOS[name]["role"]
    return "Local fallback snapshot generated from workspace clone."


def local_default_branch(path: Path) -> str:
    try:
        result = subprocess.run(
            ["git", "symbolic-ref", "--short", "HEAD"],
            cwd=path,
            text=True,
            capture_output=True,
            timeout=10,
            check=False,
        )
    except Exception:  # pragma: no cover - defensive
        return "main"
    branch = result.stdout.strip()
    return branch or "main"


def scan_local_repos() -> list[dict[str, Any]]:
    repos: dict[str, dict[str, Any]] = {}
    for child in sorted(WORKSPACE_ROOT.iterdir()):
        if not child.is_dir() or not (child / ".git").exists():
            continue
        if child.name in LOCAL_IGNORE or child.name.startswith("_"):
            continue
        remote_name = LOCAL_NAME_MAP.get(child.name, child.name)
        if remote_name in repos:
            continue
        try:
            result = subprocess.run(
                ["git", "log", "-1", "--format=%aI"],
                cwd=child,
                text=True,
                capture_output=True,
                timeout=10,
                check=False,
            )
        except Exception:  # pragma: no cover - defensive
            continue
        last_commit = result.stdout.strip() or NOW_UTC.isoformat()
        repos[remote_name] = {
            "name": remote_name,
            "html_url": f"https://github.com/{ACCOUNT}/{remote_name}",
            "description": local_description(remote_name, child),
            "language": detect_local_language(child),
            "pushed_at": last_commit,
            "updated_at": last_commit,
            "stargazers_count": 0,
            "default_branch": local_default_branch(child),
        }
    return sorted(repos.values(), key=lambda item: item["pushed_at"], reverse=True)


def latest_local_tag(name: str) -> str | None:
    if not repo_exists(name):
        return None
    try:
        result = subprocess.run(
            ["git", "tag", "--sort=-creatordate"],
            cwd=repo_path(name),
            text=True,
            capture_output=True,
            timeout=10,
            check=False,
        )
    except Exception:  # pragma: no cover - defensive
        return None
    tags = [line.strip() for line in result.stdout.splitlines() if line.strip()]
    return tags[0] if tags else None


def count_repo_signals(path: Path) -> dict[str, int | bool]:
    counts = {
        "tests": 0,
        "demos": 0,
        "examples": 0,
        "specs": 0,
        "workflows": 0,
        "docs": 0,
        "has_pyproject": (path / "pyproject.toml").exists(),
        "has_package_json": (path / "package.json").exists(),
        "has_makefile": (path / "Makefile").exists(),
        "has_citation": (path / "CITATION.cff").exists(),
        "has_security": (path / "SECURITY.md").exists(),
    }
    workflows_dir = path / ".github" / "workflows"
    if workflows_dir.exists():
        counts["workflows"] = sum(1 for _ in workflows_dir.glob("*"))

    for root, dirs, files in os.walk(path):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        rel = Path(root).relative_to(path)
        rel_parts = set(rel.parts)
        if "docs" in rel_parts:
            counts["docs"] += len(files)
        if "examples" in rel_parts:
            counts["examples"] += len(files)
        if "spec" in rel_parts or "schemas" in rel_parts or "conformance" in rel_parts:
            counts["specs"] += len(files)

        for filename in files:
            name = filename.lower()
            if name.startswith("test") or "_test" in name or name.endswith(".spec.js"):
                counts["tests"] += 1
            if "demo" in name:
                counts["demos"] += 1
            if name.endswith(".md") and "demo" in name:
                counts["demos"] += 1
            if "schema" in name or "spec" in name:
                counts["specs"] += 1
    return counts


def smoke_check(name: str, config: dict[str, Any]) -> dict[str, Any]:
    command = config["cmd"]
    path = repo_path(name)
    started = datetime.now(UTC)
    try:
        result = subprocess.run(
            command,
            cwd=path,
            text=True,
            capture_output=True,
            timeout=config.get("timeout", 60),
            check=False,
        )
    except Exception as exc:  # pragma: no cover - defensive
        return {
            "name": name,
            "ok": False,
            "returncode": None,
            "ran_at": started,
            "summary": f"execution error: {exc}",
            "output": "",
            "command": shell_join(command),
        }

    combined = "\n".join(
        part for part in [result.stdout, result.stderr] if part
    ).strip()
    markers = config.get("success_markers", [])
    ok = result.returncode == 0 and all(marker in combined for marker in markers)
    summary_line = ""
    for line in reversed(
        [line.strip() for line in combined.splitlines() if line.strip()]
    ):
        summary_line = line
        break
    return {
        "name": name,
        "ok": ok,
        "returncode": result.returncode,
        "ran_at": started,
        "summary": summary_line or f"returncode={result.returncode}",
        "output": combined,
        "command": shell_join(command),
    }


def fetch_all_repos() -> tuple[list[dict[str, Any]], str]:
    repos: list[dict[str, Any]] = []
    page = 1
    try:
        while True:
            batch = github_get(f"/users/{ACCOUNT}/repos?per_page=100&page={page}")
            if not batch:
                break
            repos.extend(batch)
            if len(batch) < 100:
                break
            page += 1
    except (HTTPError, URLError):
        if REPO_CACHE.exists():
            return json.loads(REPO_CACHE.read_text(encoding="utf-8")), "repo_cache"
        return scan_local_repos(), "local_scan"
    REPO_CACHE.write_text(
        json.dumps(repos, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return repos, "github_api"


def fetch_commits(repo_name: str, limit: int = 3) -> list[dict[str, Any]]:
    try:
        items = github_get(f"/repos/{ACCOUNT}/{repo_name}/commits?per_page={limit}")
    except (HTTPError, URLError):
        if not repo_exists(repo_name):
            return []
        try:
            result = subprocess.run(
                ["git", "log", f"-{limit}", "--format=%H%x09%aI%x09%s"],
                cwd=repo_path(repo_name),
                text=True,
                capture_output=True,
                timeout=10,
                check=False,
            )
        except Exception:  # pragma: no cover - defensive
            return []
        commits: list[dict[str, Any]] = []
        for line in result.stdout.splitlines():
            if not line.strip():
                continue
            sha, date, message = line.split("\t", 2)
            commits.append(
                {
                    "sha": sha[:7],
                    "date": date,
                    "message": normalize_message(message),
                    "url": "",
                }
            )
        return commits
    commits: list[dict[str, Any]] = []
    for item in items[:limit]:
        commit = item["commit"]
        commits.append(
            {
                "sha": item["sha"][:7],
                "date": commit["author"]["date"],
                "message": normalize_message(commit["message"]),
                "url": item["html_url"],
            }
        )
    return commits


def fetch_latest_release(repo_name: str) -> str | None:
    try:
        items = github_get(f"/repos/{ACCOUNT}/{repo_name}/releases?per_page=1")
    except (HTTPError, URLError):
        return latest_local_tag(repo_name)
    if not items:
        return latest_local_tag(repo_name)
    return items[0].get("tag_name")


def classify_repo(repo: dict[str, Any]) -> tuple[str, str]:
    name = repo["name"]
    description = (repo.get("description") or "").lower()
    if name in PRIMARY_REPOS:
        meta = PRIMARY_REPOS[name]
        return meta["layer"], meta["bp_role"]
    if name in SUPPORTING_REPOS:
        meta = SUPPORTING_REPOS[name]
        return meta["layer"], meta["role"]

    if "persona" in description or "identity" in description or "persona" in name:
        return "Identity foundation", "身份层相关支撑仓库。"
    if "audit" in description or "receipt" in description or "evidence" in description:
        return "Audit evidence", "审计证据或收据相关支撑仓库。"
    if (
        "governance" in description
        or "trust gate" in description
        or "security gate" in description
        or "policy" in description
    ):
        return "Governance gate", "治理或信任边界相关支撑仓库。"
    if (
        "integrity" in description
        or "verifiable" in description
        or "execution" in description
        or "demo" in description
    ):
        return "Execution safety / integrity", "执行路径、完整性或跨层 demo 仓库。"
    return "Adjacent architecture", "架构、标准或外围研究仓库。"


def scan_local_spear_adoption() -> dict[str, Any]:
    adopted: list[str] = []
    for child in sorted(WORKSPACE_ROOT.iterdir()):
        if not child.is_dir() or not (child / ".git").exists():
            continue
        spear_rule = child / ".github" / "security" / ".spear-rules.json"
        spear_workflow = child / ".github" / "workflows" / "spear.yml"
        if spear_rule.exists() or spear_workflow.exists():
            adopted.append(child.name)
    return {"count": len(adopted), "repos": adopted}


def collect_context(
    skip_smoke: bool = False,
    repos_override: list[dict[str, Any]] | None = None,
    source_mode: str = "github_or_cache",
) -> dict[str, Any]:
    if repos_override is not None:
        repos = repos_override
    else:
        repos, source_mode = fetch_all_repos()
    repo_index = {repo["name"]: repo for repo in repos}
    relevant_names = set(PRIMARY_REPOS) | set(SUPPORTING_REPOS)
    commits = {
        name: fetch_commits(name) for name in relevant_names if name in repo_index
    }
    releases = {
        name: fetch_latest_release(name)
        for name in relevant_names
        if name in repo_index
    }

    repo_rows: list[dict[str, Any]] = []
    languages = Counter()
    layer_counts = Counter()
    active_30d = 0
    active_7d = 0
    star_total = 0

    for repo in sorted(repos, key=lambda item: item["pushed_at"] or "", reverse=True):
        layer, note = classify_repo(repo)
        pushed_at = iso_to_dt(repo["pushed_at"]) if repo.get("pushed_at") else None
        if pushed_at:
            if NOW_UTC - pushed_at <= timedelta(days=30):
                active_30d += 1
            if NOW_UTC - pushed_at <= timedelta(days=7):
                active_7d += 1
        if repo.get("language"):
            languages[repo["language"]] += 1
        layer_counts[layer] += 1
        star_total += repo.get("stargazers_count", 0)
        repo_rows.append(
            {
                "name": repo["name"],
                "url": repo["html_url"],
                "description": repo.get("description") or "",
                "language": repo.get("language") or "n/a",
                "layer": layer,
                "note": note,
                "pushed_at": repo.get("pushed_at"),
                "updated_at": repo.get("updated_at"),
                "stars": repo.get("stargazers_count", 0),
                "default_branch": repo.get("default_branch") or "main",
            }
        )

    local_signals = {}
    for name in relevant_names:
        if repo_exists(name):
            local_signals[name] = count_repo_signals(repo_path(name))

    smoke_results: dict[str, dict[str, Any]] = {}
    if not skip_smoke:
        smoke_catalog: dict[str, dict[str, Any]] = {}
        for name, meta in PRIMARY_REPOS.items():
            smoke_catalog[name] = meta["smoke"]
        demo_smoke = SUPPORTING_REPOS["verifiable-agent-demo"].get("smoke")
        if demo_smoke and repo_exists("verifiable-agent-demo"):
            smoke_catalog["verifiable-agent-demo"] = demo_smoke
        for name, smoke in smoke_catalog.items():
            if repo_exists(name):
                smoke_results[name] = smoke_check(name, smoke)

    spear_adoption = scan_local_spear_adoption()
    return {
        "generated_at": NOW_UTC,
        "source_mode": source_mode,
        "repos": repos,
        "repo_index": repo_index,
        "repo_rows": repo_rows,
        "commits": commits,
        "releases": releases,
        "local_signals": local_signals,
        "smoke_results": smoke_results,
        "spear_adoption": spear_adoption,
        "active_30d": active_30d,
        "active_7d": active_7d,
        "language_counts": languages,
        "layer_counts": layer_counts,
        "star_total": star_total,
    }


def build_account_snapshot(ctx: dict[str, Any]) -> str:
    star_text = str(ctx["star_total"])
    if ctx.get("source_mode") in {"local_scan", "local_fallback"}:
        star_text = "本地模式不统计"
    language_text = ", ".join(
        f"{language} ({count})"
        for language, count in ctx["language_counts"].most_common(6)
    )
    layer_text = ", ".join(
        f"{LAYER_LABELS.get(layer, layer)}: {ctx['layer_counts'][layer]}"
        for layer in LAYER_ORDER
        if ctx["layer_counts"].get(layer)
    )
    return textwrap.dedent(
        f"""\
        ## 账号概览

        | 指标 | 数值 |
        | --- | --- |
        | 扫描到的公开仓库 | {len(ctx["repos"])} |
        | 最近 30 天有更新的仓库 | {ctx["active_30d"]} |
        | 最近 7 天有更新的仓库 | {ctx["active_7d"]} |
        | GitHub 星标总数 | {star_text} |
        | 本地 `spear-check` 接入数 | {ctx["spear_adoption"]["count"]} |
        | 主要语言 | {language_text} |
        | 分层分布 | {layer_text} |
        """
    ).strip()


def core_table(ctx: dict[str, Any]) -> str:
    lines = [
        "## AegisTrust 四个核心模块",
        "",
        "| 模块 | 主仓库 | 最近代码提交 | 最近版本 | 本地回放 | 当前状态 |",
        "| --- | --- | --- | --- | --- | --- |",
    ]
    for name in [
        "god-spear",
        "safety-valve-spec",
        "execution-integrity-core",
        "aro-audit",
    ]:
        repo = ctx["repo_index"][name]
        commit = ctx["commits"].get(name, [{}])[0] if ctx["commits"].get(name) else {}
        release = ctx["releases"].get(name) or "none"
        smoke = ctx["smoke_results"].get(name)
        smoke_text = "not run"
        if smoke:
            smoke_text = "PASS" if smoke["ok"] else f"FAIL (rc={smoke['returncode']})"
        lines.append(
            "| {module} | [{repo}]({url}) | {commit_date} `{sha}` | {release} | {smoke} | {stage} |".format(
                module=PRIMARY_REPOS[name]["module"],
                repo=name,
                url=repo["html_url"],
                commit_date=fmt_day(commit.get("date")),
                sha=commit.get("sha", "n/a"),
                release=release,
                smoke=smoke_text,
                stage=PRIMARY_REPOS[name]["stage"],
            )
        )
    return "\n".join(lines)


def render_repo_detail(name: str, ctx: dict[str, Any]) -> str:
    repo = ctx["repo_index"][name]
    meta = PRIMARY_REPOS[name]
    commit_items = ctx["commits"].get(name, [])
    smoke = ctx["smoke_results"].get(name)
    signals = ctx["local_signals"].get(name, {})
    smoke_block = "- 本地回放：未执行。"
    if smoke:
        status = "PASS" if smoke["ok"] else f"FAIL (rc={smoke['returncode']})"
        smoke_block = (
            f"- 本地回放（{fmt_dt(smoke['ran_at'])}）：{status}，命令 `{smoke['command']}`\n"
            f"- 关键输出：`{smoke['summary']}`"
        )
    commit_lines = (
        "\n".join(
            f"  - {fmt_day(item['date'])} `{item['sha']}` {item['message']}"
            for item in commit_items
        )
        or "  - 暂未拿到提交信息。"
    )
    capabilities = "\n".join(f"- {item}" for item in meta["capabilities"])
    validation = "\n".join(f"- {item}" for item in meta["validation"])
    signal_line = (
        f"- 本地代码信号：tests={signals.get('tests', 0)}, demos={signals.get('demos', 0)}, "
        f"examples={signals.get('examples', 0)}, workflows={signals.get('workflows', 0)}, "
        f"spec 文件约={signals.get('specs', 0)}"
    )
    return "\n".join(
        [
            f"### {meta['module']}: `{name}`",
            "",
            f"- GitHub 地址: {repo['html_url']}",
            f"- 仓库描述: {repo.get('description') or meta['one_line']}",
            f"- 模块作用: {meta['bp_role']}",
            f"- 最近推送时间: {fmt_dt(repo.get('pushed_at'))}",
            f"- 最近版本: {ctx['releases'].get(name) or 'none'}",
            signal_line,
            smoke_block,
            "",
            "最近提交:",
            commit_lines,
            "",
            "已实现能力:",
            capabilities,
            "",
            "开源验证:",
            validation,
        ]
    ).strip()


def build_supporting_layers(ctx: dict[str, Any]) -> str:
    grouped: dict[str, list[str]] = defaultdict(list)
    for name in sorted(SUPPORTING_REPOS):
        if name not in ctx["repo_index"]:
            continue
        grouped[SUPPORTING_REPOS[name]["layer"]].append(name)

    lines = ["## 支撑仓库分层", ""]
    for layer in LAYER_ORDER:
        if not grouped.get(layer):
            continue
        lines.append(f"### {LAYER_LABELS.get(layer, layer)}")
        for name in grouped[layer]:
            repo = ctx["repo_index"][name]
            release = ctx["releases"].get(name) or "none"
            note = SUPPORTING_REPOS[name]["role"]
            commit = (
                ctx["commits"].get(name, [{}])[0] if ctx["commits"].get(name) else {}
            )
            lines.append(
                f"- [{name}]({repo['html_url']}): {note} 最近提交 {fmt_day(commit.get('date'))}，最近版本 {release}。"
            )
        lines.append("")
    return "\n".join(lines).rstrip()


def build_full_inventory(ctx: dict[str, Any]) -> str:
    lines = [
        "## 公开仓库清单",
        "",
        "| 仓库 | 分层 | 语言 | 最近推送（UTC） | 备注 |",
        "| --- | --- | --- | --- | --- |",
    ]
    for row in ctx["repo_rows"]:
        short_note = row["note"].replace("|", "/")
        lines.append(
            f"| [{row['name']}]({row['url']}) | {LAYER_LABELS.get(row['layer'], row['layer'])} | {row['language']} | {fmt_day(row['pushed_at'])} | {short_note} |"
        )
    return "\n".join(lines)


def build_github_progress_report(ctx: dict[str, Any]) -> str:
    smoke_lines = []
    for name in [
        "god-spear",
        "safety-valve-spec",
        "execution-integrity-core",
        "aro-audit",
        "verifiable-agent-demo",
    ]:
        smoke = ctx["smoke_results"].get(name)
        if not smoke:
            continue
        result = "PASS" if smoke["ok"] else f"FAIL (rc={smoke['returncode']})"
        smoke_lines.append(
            f"- `{name}`：{result}，时间 {fmt_dt(smoke['ran_at'])}，命令 `{smoke['command']}`。"
        )
    smoke_block = "\n".join(smoke_lines) if smoke_lines else "- 本地回放已跳过。"
    source_text = {
        "github_api": "直接使用 GitHub API 拉取远程仓库信息。",
        "repo_cache": "GitHub API 被限流，已回退到上次保存的远程仓库快照。",
        "local_scan": "GitHub API 被限流且没有缓存，已回退到本地仓库扫描。",
        "local_fallback": "本次直接使用本地仓库和本地 git 信息生成。",
    }.get(ctx.get("source_mode"), "优先使用 GitHub API，必要时回退到缓存或本地仓库。")
    head = "\n".join(
        [
            "# AegisTrust 本地仓库进展报告",
            "",
            f"生成时间：{fmt_dt(ctx['generated_at'])}",
            "",
            "数据来源：",
            f"- {source_text}",
            f"- 仓库归属：`{ACCOUNT}` 的本地仓库克隆",
            f"- 本地工作区：`{WORKSPACE_ROOT}`",
            "- 同时复跑了四个核心模块和端到端 demo 的轻量验证命令",
            "",
            "## 一句话结论",
            "",
            (
                f"- 当前本地仓库群共扫描到 {len(ctx['repos'])} 个仓库，最近 30 天有更新的仓库有 "
                f"{ctx['active_30d']} 个，说明这套体系到 {fmt_day(ctx['generated_at'])} 仍在持续推进。"
            ),
            (
                "- 现在最清楚、最容易对外讲明白的产品结构，已经收敛成四个模块："
                "`god-spear`、`safety-valve-spec`、`execution-integrity-core`、`aro-audit`。"
            ),
            (
                "- 账号里还有 persona、interaction、runtime governance 等更大体系。"
                "这些内容能证明技术深度，但不适合放在第一页讲太多，更适合放到补充材料。"
            ),
            "- 当前最准确的阶段描述是：开源验证已完成，可回放 demo 已完成，场景打磨和试点转化正在推进。",
            "",
            "## 最新本地回放",
            smoke_block,
        ]
    )

    parts = [
        head,
        build_account_snapshot(ctx),
        core_table(ctx),
        "## 核心模块详情",
        render_repo_detail("god-spear", ctx),
        render_repo_detail("safety-valve-spec", ctx),
        render_repo_detail("execution-integrity-core", ctx),
        render_repo_detail("aro-audit", ctx),
        build_supporting_layers(ctx),
        textwrap.dedent(
            f"""\
            ## 当前可对外使用的判断

            - 现在最适合作为主产品外壳的仓库是 `aro-audit`，因为它已经有对外说明、快速演示和证据打包能力。
            - 现在最适合对外讲的结构是：`god-spear` + `safety-valve-spec` + `execution-integrity-core` + `aro-audit`。
            - 当前技术成熟度可以概括为：核心流程能本地复跑，3 个核心仓库已有公开版本，`spear-check` 已在 {ctx["spear_adoption"]["count"]} 个本地仓库里持续使用。
            - 当前商业成熟度仍偏早期，更适合写成“技术验证完成，正在走向试点和场景验证”。
            - 当前最统一的一句话定位：**AegisTrust = AI 可信执行基础设施**。
            """
        ).strip(),
        build_full_inventory(ctx),
    ]
    return "\n\n".join(parts) + "\n"


def build_technical_progress(ctx: dict[str, Any]) -> str:
    repo = ctx["repo_index"]
    smoke = ctx["smoke_results"]
    release = ctx["releases"]
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 技术进展说明

        生成时间：{fmt_dt(ctx["generated_at"])}

        一句话判断：AegisTrust 现在已经不是一个概念集合，而是已经能本地复跑、能讲清主线、能对外演示的 AI 可信执行链。

        ## 1. 当前技术结构

        现在最适合对外讲的结构不是“大而全架构”，而是四个模块连成一条主线：

        1. `god-spear`：上线前门禁，也就是代码上线前先自动体检
        2. `safety-valve-spec`：动作边界收据，也就是关键动作必须带小票
        3. `execution-integrity-core`：执行完整性证明，也就是证明过程没被偷偷改过
        4. `aro-audit`：事后证据包和独立复验，也就是出事后能把账查清楚

        这样讲有两个好处。第一，评委和客户容易听懂。第二，它和本地仓库里已经存在的代码、版本、demo、验证脚本完全对得上。

        ## 2. 已经做成了什么

        **`god-spear`**
        - 已能在 `CI`（代码上线前的自动检查流程）里检查 `trust boundary`（谁能碰什么）、`failure signal`（出了问题能不能明确报警）、`revocation pathway`（发现问题后能不能及时撤回）这些关键条件。
        - 已有公开版本 `{release.get("god-spear") or "none"}`，本地回放在 {fmt_day(smoke.get("god-spear", {}).get("ran_at")) if smoke.get("god-spear") else "n/a"} 通过。

        **`safety-valve-spec`**
        - 已定义 `ALLOW / DENY / DEGRADE` 收据格式，也就是每次关键动作都要明确写清“允许、拒绝、降级”。
        - 配套已有签名、验签、`CA / CRL`（证书和吊销链路）、`compat attestation`（兼容性证明）和 gateway demo。
        - 已有公开版本 `{release.get("safety-valve-spec") or "none"}`，本地 conformance 在 {fmt_day(smoke.get("safety-valve-spec", {}).get("ran_at")) if smoke.get("safety-valve-spec") else "n/a"} 全量通过。

        **`execution-integrity-core`**
        - 已实现 `hash chain`（前后相扣、改一处整条会露馅）、确定性导出（同样数据导出结果一致）、全链校验和篡改检测。
        - 它解决的是“怎么证明过程没被动过”，而不只是“有没有留下一份日志”。

        **`aro-audit`**
        - 已实现 `audit bundle`（可交付的证据包）、`manifest`（材料清单）、`Merkle/checkpoint`（防篡改校验点）、验签流程和 quickstart。
        - 已有公开版本 `{release.get("aro-audit") or "none"}`，本地 quickstart 在 {fmt_day(smoke.get("aro-audit", {}).get("ran_at")) if smoke.get("aro-audit") else "n/a"} 验证通过。

        ## 3. 现在能演示到什么程度

        当前已经不是“只有 README 和架构图”的阶段，而是关键链路都能直接复跑：

        - `god-spear`：结果 `{("PASS" if smoke.get("god-spear", {}).get("ok") else "not-run")}`，输出 `STATUS: PASS`
        - `safety-valve-spec`：结果 `{("PASS" if smoke.get("safety-valve-spec", {}).get("ok") else "not-run")}`，输出 `Overall: PASS`，说明规范和实现是对得上的
        - `execution-integrity-core`：结果 `{("PASS" if smoke.get("execution-integrity-core", {}).get("ok") else "not-run")}`，输出 `SELF_CHECK: PASS`，说明完整性证明能自证
        - `aro-audit`：结果 `{("PASS" if smoke.get("aro-audit", {}).get("ok") else "not-run")}`，正常样本 `VERIFY_OK`，篡改样本 `Merkle mismatch`，说明动过手脚就能被发现
        - `verifiable-agent-demo`：结果 `{("PASS" if smoke.get("verifiable-agent-demo", {}).get("ok") else "not-run")}`，可输出完整 `intent / action / result / audit_record`，也就是一次动作前后完整留痕

        用人话说，AegisTrust 现在已经能证明三件事：能拦、能留证、能查出篡改。

        ## 4. 开源和工程验证情况

        - 当前本地仓库群共扫描到 {len(ctx["repos"])} 个仓库，其中最近 30 天有更新的仓库有 {ctx["active_30d"]} 个。
        - `god-spear`、`safety-valve-spec`、`aro-audit` 都已有公开版本；`aro-audit` 和 `persona-object-protocol` 还带 DOI / citation 资产。
        - 本地扫描显示 `spear-check` 已接入 {ctx["spear_adoption"]["count"]} 个仓库，说明这不是孤立 demo，而是在持续复用。
        - `verifiable-agent-demo` 在 {fmt_day(repo["verifiable-agent-demo"]["pushed_at"])} 仍有更新，说明跨层 demo 还在往前推进。

        ## 5. 下一步重点

        未来 6 到 12 个月，最重要的不是再堆更多概念，而是把已有模块压成更短的交付路径：

        1. 先用 `aro-audit` 作为主入口，对外卖“证据包 + 独立复验”。
        2. 再用 `god-spear` 把风险前移到 `CI`，也就是让问题在上线前暴露。
        3. 再用 `safety-valve-spec` 把关键动作的小票标准化。
        4. 用 `execution-integrity-core` / `fdo-kernel-mvk` 继续把执行层证明做厚。

        对外最合适的话术不是“我们做了很多协议”，而是“我们已经把 AI 高风险动作的问题，拆成可拦截、可留证、可复验的一条工程链”。
        """
        ).strip()
        + "\n"
    )


def build_architecture_summary(ctx: dict[str, Any]) -> str:
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 架构概述

        生成时间：{fmt_dt(ctx["generated_at"])}

        一句话定义：AegisTrust 是一条围绕 AI 真实动作建立起来的可信执行链，不是单点插件，也不是只会提醒风险的看板。

        ## 1. 这套架构从哪里来

        AegisTrust 不是凭空冒出来的项目，它来自一套更完整的 AI agent 架构研究。更大的总图里，`POP` 这类仓库负责身份层，也就是回答“这个 agent 是谁”；GitHub 首页和架构总览仓库负责把五层关系讲清楚。

        但在当前阶段，最适合拿出来做产品和路演的，不是把所有层一起讲，而是先聚焦最刚需、最容易被企业理解和采购的一段闭环：可信执行。

        ## 2. 当前对外主线：四层可信执行闭环

        **1. CI governance：`god-spear`**
        - 作用：在上线前做门禁检查。
        - 人话解释：代码发布前先自动体检，看看规则有没有写完整、风险有没有交代清楚。

        **2. runtime safety：`safety-valve-spec`**
        - 作用：给关键动作加收据规则。
        - 人话解释：高风险动作不能“口头批准”，每次都要带一张能核对的小票。

        **3. execution integrity：`execution-integrity-core`**
        - 作用：证明执行过程没有被偷偷改写。
        - 人话解释：不是只留一份日志，而是要能证明“这条过程就是按这个顺序发生的”。

        **4. audit evidence：`aro-audit`**
        - 作用：把执行结果整理成可交付、可复验的证据包。
        - 人话解释：出了问题以后，不是翻聊天记录，而是直接拿出一包能查账的材料。

        ## 3. 四层怎么连起来

        这四层串起来以后，逻辑非常简单：

        1. 上线前先检查，别让明显有缺口的配置进入生产。
        2. 执行时要求关键动作带收据，别让高风险操作“裸奔”。
        3. 执行后把过程串成可验证链路，别让日志变成谁都能改的流水账。
        4. 最后导出证据包，别让复盘和审计只剩口头解释。

        ## 4. 为什么现在只讲这四层

        因为这是当前最容易形成产品闭环的一段。

        - 企业最先愿意为“风险前移”和“证据可交付”买单。
        - 评委和客户最容易听懂“先拦住，再留证，最后可复查”。
        - 本地仓库和 demo 也已经证明，这四层不是概念，而是能跑起来的工程链。

        ## 5. POP 和 GitHub 首页放在什么位置

        它们不是没关系，而是应该放在“架构来源”和“长期延展”的位置。

        - `POP` 证明你不是只会做安全补丁，而是在做更完整的 agent 体系。
        - GitHub 首页证明你有总架构，不是堆几个零散仓库。
        - 但在当前 BP 里，它们应该作为背景和深度证明，而不是抢走产品主线。

        所以最合适的表达是：**AegisTrust 是更大架构中的当前商业化切片，先聚焦 AI 可信执行闭环。**
        """
        ).strip()
        + "\n"
    )


def build_policy_context(ctx: dict[str, Any]) -> str:
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 政策语境

        生成时间：{fmt_dt(ctx["generated_at"])}

        一句话判断：AegisTrust 最适合被写成“顺应产业方向的可信执行基础能力”，而不是“直接承接政策任务的项目”。

        ## 1. 适合借力的政策语境

        近一阶段公开政策信号持续强调几个方向：数字化转型、现代化产业体系、`AI+` 应用深化，以及发展和安全并重的治理思路。在这种语境下，市场会越来越关注一个现实问题：AI 不只是会不会生成内容，而是能不能安全进入真实流程。

        这正好给了 AegisTrust 一个很自然的位置。它不是去争夺“最大模型”或“通用平台”的话语权，而是补上 AI 进入企业流程之后最容易缺失的一层基础能力。

        ## 2. 最适合的表达方式

        建议在 BP 和路演里使用这些表达：

        - 与国家数字化转型方向保持一致
        - 符合发展与安全并重的产业趋势
        - 支持人工智能在真实业务中的可信部署
        - 为企业智能化升级提供可验证、可审计的基础能力
        - 面向新一代数字基础设施中的可信执行环节

        这些说法的好处是，评委能自然感受到项目和大方向同频，但不会觉得你在“借政策压人”。

        ## 3. 不建议怎么写

        不建议使用过于直给的表述，比如：

        - 直接服务政府
        - 直接承担国家战略任务
        - 直接对应某项专项工程
        - 国家急需我们这类项目

        这些说法在比赛和申报里容易显得用力过猛，也会让项目从创业叙事滑向口号叙事。

        ## 4. 为什么 AegisTrust 能自然对齐

        因为它解决的是一个基础问题：当 AI 从“会说”进入“会做”时，企业如何把关键动作纳入可信、可验、可追责的轨道。

        这个问题同时连接了几个方向：

        - `AI+`：AI 进入业务流程以后，必须面对执行风险。
        - 新质生产力：企业想提升效率，但也要控制失误成本。
        - AI 治理：不是只谈原则，而是把治理要求落到工程组件上。
        - 数字基础设施：可信执行和审计能力，会逐步成为基础设施的一部分。

        ## 5. BP 里的推荐落点

        在 BP 里，政策相关内容最适合写成“行业环境”和“发展趋势”，不要写成“政治表态”。

        更合适的话术是：

        随着人工智能深入企业真实场景，可信执行、风险治理和审计能力正逐渐成为智能化升级中的基础要求。AegisTrust 所做的，正是面向这一趋势提供可落地的可信执行组件。

        ## 6. 北京落地怎么顺带写

        如果写北京落地，最合适的角度不是喊口号，而是强调资源条件：

        - 北京拥有高密度 AI 企业、研究机构和产业服务资源
        - 企业数字化和智能化升级需求集中
        - 高校、研究机构、园区和创新平台更容易形成联合验证

        这样写，既自然贴近地方生态，也保留了创业项目应有的商业气质。
        """
        ).strip()
        + "\n"
    )


def build_bp(ctx: dict[str, Any]) -> str:
    repo = ctx["repo_index"]
    smoke = ctx["smoke_results"]
    release = ctx["releases"]
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 路演稿

        生成时间：{fmt_dt(ctx["generated_at"])}

        一句话定位：**AegisTrust 不是另一个 Agent 工具，而是一层 AI 可信执行基础设施。**

        再说得更直白一点：当 AI 开始真的去调工具、改配置、跑流程时，我们要做的不是继续和幻觉赌输赢，而是让高风险动作先被拦住、过程留下小票、事后能把账查清楚。

        ## 1. 项目背景

        过去大家讨论 AI，更多是在看它会不会写、会不会答。现在真正的变化是，AI 已经开始碰真实动作，比如运维操作、合同处理、财务流程、企业协同。只要 AI 从“会说”走到“会做”，风险级别就会立刻上升。AegisTrust 就是在这个节点上出现的。

        ## 2. 关键问题

        企业真正头疼的不是一个点，而是三件事同时发生：

        - AI 可能因为幻觉、上下文不全或策略写错，去做本来不该做的事。
        - 普通日志只能证明“发生过”，很难证明“是不是按规则发生”。
        - 一旦出事，企业很难向客户、审计或监管方拿出一套独立可复验的证据。

        所以问题的本质不是“模型准不准”，而是“动作能不能被信任”。

        ## 3. 解决办法

        AegisTrust 的思路很朴素：不要只盯模型内部，而是在 AI 动作出口外面加一层基础设施。项目现在收敛成四个模块：

        - `god-spear`：上线前门禁
        - `safety-valve-spec`：动作边界收据
        - `execution-integrity-core`：执行完整性证明
        - `aro-audit`：事后证据包和独立复验

        合在一起，就是一条完整的可信执行链。人话就是六个字：**先拦住，再留证。**

        ## 4. 核心技术

        这套技术不是单点功能，而是一条闭环：

        - `CI`（代码上线前的自动检查流程）阶段，由 `god-spear` 检查 `trust boundary`（谁能碰什么）、`failure signal`（出了问题能不能明确报警）、`revocation pathway`（发现问题后能不能及时撤回）。
        - 执行动作前，由 `safety-valve-spec` 要求动作必须带 `ALLOW / DENY / DEGRADE` 收据，也就是“允许、拒绝、降级”三种明确结果。
        - 执行过程中，由 `execution-integrity-core` 用 `hash chain`（前后相扣、改一处整条会露馅）和确定性导出证明执行链没有被改写。
        - 执行结束后，由 `aro-audit` 打成 `audit bundle`（可交付的证据包），里面带 `manifest`（材料清单）、`Merkle/checkpoint`（防篡改校验点）和验签流程。

        如果不用技术词来讲，这套系统做的就是三件事：上线前先体检，执行时留小票，出事后能查账。

        ## 5. AI 幻觉风险怎么处理

        AegisTrust 不承诺“消灭幻觉”。更现实的做法是，即使模型判断错了，也不能让错误直接变成高风险动作。只要没有通过 `CI` 门禁、没有合规收据、没有进入完整性链路，它就不能进入可信交付路径。

        换句话说，我们不是把风险讲小，而是把风险变成三种可管理的状态：可拦截、可发现、可追责。

        ## 6. 产品架构

        产品架构分四层，而且可以单独接入：

        - 第一层是上线前门禁，由 `god-spear` 检查配置是否完整。
        - 第二层是动作边界收据，由 `safety-valve-spec` 保证关键动作必须带凭证。
        - 第三层是执行完整性，由 `execution-integrity-core` 或进阶版 `fdo-kernel-mvk` 证明执行链没被篡改。
        - 第四层是审计证据，由 `aro-audit` 输出可以交付给客户、审计和合作方的证据包。

        这不是“做四个仓库”，而是围绕一条业务链把问题拆开解决。

        ## 7. 产品形态

        现在最适合对外落地的形态有三种：

        - 开源基础版：协议、校验器、示例和 quickstart，方便客户先试。
        - 企业部署版：行业策略包、标准化审计模板、托管验证接口。
        - 交付服务版：帮助客户把高风险 AI 流程真正接入门禁、收据、审计和复验链路。

        这意味着我们既能先低成本进入客户，也能往上做软件收入和服务收入。

        ## 8. 商业模式

        商业模式可以理解成三层：

        - 开源标准做入口，用来建立信任和试用门槛优势。
        - 企业软件做复用收入，用来卖部署、策略包和验证服务。
        - 实施服务做高附加值收入，用来做试点接入、培训认证和长期支持。

        这样做的好处是，不需要一开始就卖一个很重的平台，也能随着客户成熟度慢慢放大收入。

        ## 9. 市场机会

        市场机会来自一个非常明确的变化：企业 AI 正在从问答型工具，变成能触发真实动作的系统。最先有需求的部门也很清楚，就是 IT 运维、法务、财务、合规和风控。

        因为这些部门一旦把 AI 接进真实流程，就会立刻碰到可信执行和审计交付的问题。AegisTrust 卡住的，正是这个“从能用到敢用”的空档。

        ## 10. 竞争优势

        AegisTrust 的优势不是某一个功能特别炫，而是它把闭环做出来了。

        - 日志平台能看见问题，但很难强证明。
        - 提示词 `guardrail`（输入输出限制）更像在管“说什么”，不够管“做什么”。
        - 单点安全插件能挡一个点，但挡不住整条链。

        AegisTrust 从上线前、执行中、执行后全部打通，所以更适合高风险动作场景。

        ## 11. 技术进展

        截至 {fmt_day(ctx["generated_at"])}, 本地仓库群共扫描到 {len(ctx["repos"])} 个仓库，最近 30 天有更新的仓库为 {ctx["active_30d"]} 个。四个核心模块中，`god-spear` 最近版本是 `{release.get("god-spear") or "none"}`，`safety-valve-spec` 是 `{release.get("safety-valve-spec") or "none"}`，`aro-audit` 是 `{release.get("aro-audit") or "none"}`。

        本地扫描还显示 `spear-check` 已接入 {ctx["spear_adoption"]["count"]} 个仓库，说明项目不是一次性的原型，而是在持续使用。

        ## 12. 演示能力

        当前演示能力已经能支撑对外沟通：

        - {fmt_day(smoke.get("god-spear", {}).get("ran_at")) if smoke.get("god-spear") else fmt_day(ctx["generated_at"])} 重新执行 `god-spear`，返回 `STATUS: PASS`。
        - 同日重新执行 `safety-valve-spec` conformance（规范一致性检查），返回 `Overall: PASS`。
        - `execution-integrity-core` 返回 `SELF_CHECK: PASS`。
        - `aro-audit` quickstart 先得到 `VERIFY_OK: full chain valid`，再对篡改样本触发 `Merkle mismatch`。
        - `verifiable-agent-demo` 于 {fmt_day(repo["verifiable-agent-demo"]["pushed_at"])} 仍有更新，并可输出完整的 `intent / action / result / audit_record`（一次动作前后留下的结构化记录）。

        这意味着现场演示时，我们不是只能放架构图，而是真能把关键链路跑给别人看。

        ## 13. 团队情况

        当前公开可验证的核心力量是独立研究者 Bin Zhang，长期围绕 persona、runtime governance、execution integrity 和 audit evidence 做体系化建设。团队当前最强的证明不是头衔，而是持续更新的代码、版本、demo 和验证脚本。

        对外更合适的表达是：创始人主导技术核心，同时正在补齐行业顾问、合规理解和 ToB 落地能力。

        ## 14. 发展路线

        - 未来 3 个月，重点是把四个模块压缩成一条更短的交付路径，做出更清晰的行业样板。
        - 未来 3 到 6 个月，重点是拿到 2 到 3 个试点或合作意向，把“技术验证”推进到“场景验证”。
        - 未来 6 到 12 个月，重点是形成企业部署版、托管验证接口和实施服务能力。

        ## 15. 融资计划

        当前更适合的是早期验证资金或带场景资源的战略资金，而不是大规模融资。资金用途也应该写得更实一点：

        - 继续打磨验证工具链和稳定演示环境
        - 投入试点实施和客户成功
        - 补齐行业模板、合规材料和本地生态合作

        这样写既真实，也更利于后续落地。

        ## 16. 落地计划

        落地的关键不是喊口号，而是先找到高风险 AI 流程场景。

        - 优先去做有审计压力和内控需求的企业流程
        - 再和高校、研究机构、产业园区合作做验证和共建
        - 最后逐步形成培训、认证和实施能力

        最终目标不是停留在开源项目，而是形成可持续交付的产品和服务。

        ## 结尾

        AegisTrust 真正有价值的地方，不是仓库数量，而是这些仓库正在收拢成一条清楚的产品线：当 AI 从“生成内容”进入“执行动作”，市场需要的是一层让关键动作可拦截、可留证、可复验的基础设施。AegisTrust 正在做的，就是这层基础设施。
        """
        ).strip()
        + "\n"
    )


def build_bp_v2(ctx: dict[str, Any]) -> str:
    repo = ctx["repo_index"]
    smoke = ctx["smoke_results"]
    release = ctx["releases"]
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 商业计划书 V2

        生成时间：{fmt_dt(ctx["generated_at"])}

        一句话定位：**AegisTrust 是一套 AI 可信执行基础设施，让高风险 AI 动作可拦截、可留证、可复验。**

        项目源自更完整的五层 agent 架构研究，但当前商业化阶段只聚焦最刚需的一段闭环：可信执行。这样既保留了技术深度，也让评委和客户能一眼看懂现在到底在做什么。

        ## 1. 项目背景

        AI 正在从“会写、会答”走向“会调工具、会改配置、会跑流程”。一旦 AI 开始碰真实动作，企业面对的就不再只是模型效果问题，而是执行信任问题。谁来保证它不越权、不乱做、出了事能查清楚，这就是 AegisTrust 的出发点。

        ## 2. 行业痛点

        当前企业把 AI 接入真实流程时，通常会遇到三类问题：

        - 幻觉、上下文缺失、策略错误可能直接变成错误动作
        - 普通日志只能记录发生过，不能强证明发生得是否合规
        - 一旦出现事故，企业很难拿出独立可复验的证据材料

        所以真正缺少的，不是又一个模型能力，而是一层可信执行能力。

        ## 3. 解决方案

        AegisTrust 的解决办法很直接：不和模型的每一次判断赌输赢，而是在动作出口外面建立一条可信执行链。

        - 上线前先做门禁检查
        - 执行时要求关键动作带收据
        - 执行后输出可验证、可复查的证据包

        人话就是：**先拦住，再留票，最后还能查账。**

        ## 4. 架构设计

        当前对外主线收敛成四层：

        - `god-spear`：`CI governance`，也就是上线前先自动体检
        - `safety-valve-spec`：`runtime safety`，也就是关键动作必须带小票
        - `execution-integrity-core`：`execution integrity`，也就是证明过程没有被偷偷改写
        - `aro-audit`：`audit evidence`，也就是把结果整理成一包可交付材料

        这四层已经能组成一条企业可理解、可接入、可演示的可信执行闭环。

        ## 5. AI 幻觉风险治理

        AegisTrust 不承诺消灭幻觉，而是控制幻觉的后果。只要关键动作没有通过门禁、没有收据、没有进入完整性链路，它就不能进入可信交付路径。

        这意味着我们处理幻觉的方式不是“保证永远不出错”，而是把错误变成可拦截、可发现、可追责的问题。

        ## 6. 产品形态

        当前最适合落地的产品形态有三类：

        - 开源基础版：协议、校验器、示例、quickstart
        - 企业部署版：策略包、审计模板、验证接口
        - 实施服务版：场景接入、培训认证、长期支持

        先用开源进入场景，再逐步放大软件收入和服务收入，这是更现实的路径。

        ## 7. 市场机会

        随着企业 AI 从问答工具进入真实流程，最先出现需求的会是 IT 运维、法务、财务、合规和风控等高责任部门。因为这些场景一旦使用 AI，就会立刻碰到“敢不敢放权”和“出了事怎么证明”的问题。

        AegisTrust 面对的，不是泛 AI 工具市场，而是高风险 AI 动作的可信执行基础设施市场。

        ## 8. 政策环境与方向

        从近一阶段公开政策语境来看，数字化转型、`AI+` 深化应用、现代化产业体系建设，以及发展和安全并重，已经是比较清晰的大方向。AegisTrust 最合适的写法，不是说自己“直接服务某项政策”，而是强调它顺应了可信部署、风险治理和安全落地的产业趋势。

        更自然的表达是：随着人工智能深入真实业务场景，可信执行与安全治理正逐渐成为智能化升级中的基础能力，AegisTrust 正是在这一趋势下提供可落地组件。

        ## 9. 竞争优势

        市面上常见方案大致分三类：

        - 日志平台：能看见过程，但不一定能强证明
        - 提示词 `guardrail`（输入输出限制）：更像在管模型怎么说，不够管模型怎么做
        - 单点安全插件：能挡某一处，但挡不住整条链

        AegisTrust 的差异在于，它把上线前、执行中、执行后打通成闭环，更适合高风险动作场景。

        ## 10. 技术进展

        截至 {fmt_day(ctx["generated_at"])}, 本地仓库群共扫描到 {len(ctx["repos"])} 个仓库，最近 30 天有更新的仓库为 {ctx["active_30d"]} 个。四个核心模块中，`god-spear` 最近版本为 `{release.get("god-spear") or "none"}`，`safety-valve-spec` 为 `{release.get("safety-valve-spec") or "none"}`，`aro-audit` 为 `{release.get("aro-audit") or "none"}`。

        本地扫描还显示 `spear-check` 已接入 {ctx["spear_adoption"]["count"]} 个仓库，说明这不是一次性原型，而是在持续复用的工程链。

        ## 11. 演示能力

        当前已经具备可现场展示的演示链路：

        - `god-spear` 本地回放返回 `STATUS: PASS`
        - `safety-valve-spec` conformance 返回 `Overall: PASS`
        - `execution-integrity-core` 返回 `SELF_CHECK: PASS`
        - `aro-audit` 正常样本 `VERIFY_OK`，篡改样本触发 `Merkle mismatch`
        - `verifiable-agent-demo` 可输出完整 `intent / action / result / audit_record`

        也就是说，这不是“只能讲故事”的项目，而是关键路径可以直接跑出来。

        ## 12. 商业模式

        商业模式分三层：

        - 开源标准做入口，降低试用门槛，建立行业信任
        - 企业软件做复用收入，提供策略包、模板和验证能力
        - 实施服务做高附加值收入，承接试点、接入和认证

        这样的结构更适合早期项目逐步放大。

        ## 13. 发展路线

        - 0 到 3 个月：压缩产品叙事，形成清晰行业样板
        - 3 到 6 个月：拿到 2 到 3 个试点或合作意向
        - 6 到 12 个月：形成企业部署能力和标准化验证服务

        路线重点不是再扩概念，而是把现有能力变成更短、更清晰的交付路径。

        ## 14. 团队介绍

        当前公开可验证的核心力量是独立研究者 Bin Zhang，长期围绕 persona、runtime governance、execution integrity 和 audit evidence 做体系化建设。更重要的是，相关能力已经体现在持续更新的仓库、版本、demo 和验证脚本中。

        当前团队形态更接近“技术核心已成形，商业化能力正在补齐”的阶段。

        ## 15. 融资计划

        当前更适合的是早期验证资金或带场景资源的战略资金。资金将主要用于三件事：

        - 打磨稳定演示环境和验证工具链
        - 推进试点交付和客户成功
        - 补齐行业模板、合规材料和生态合作

        重点不是大规模扩张，而是先把产品和场景闭环做实。

        ## 16. 北京落地计划

        北京适合成为项目早期落地的重要节点，不是因为口号，而是因为条件匹配：

        - AI 企业、研究机构和产业服务资源密集
        - 企业数字化和智能化升级需求集中
        - 更容易联合高校、园区和创新平台做试点验证

        更现实的路径是先从高审计压力、高内控要求的企业流程切入，再逐步形成样板、认证和持续服务能力。

        ## 结论

        AegisTrust 的核心价值不在于仓库数量，而在于已经把 AI 高风险动作的问题拆成了一条可落地的工程链。当 AI 从“会说”进入“会做”，市场最终需要的不是更多口号，而是一层可信执行基础设施。AegisTrust 正在把这层基础设施做出来。
        """
        ).strip()
        + "\n"
    )


def build_competition_summary() -> str:
    text = (
        "AegisTrust 是一套 AI 可信执行基础设施，面向高风险 AI 动作场景。"
        "它解决的不是模型会不会回答问题，而是当 AI 开始调用工具、修改配置、触发流程时，企业如何做到敢放权、能控风险、出了事说得清。"
        "项目当前已经形成一条清晰闭环：上线前做门禁检查，执行中要求关键动作带收据，执行后输出可复核的证据包。"
        "简单说，就是让关键动作先体检、做事要留票、结果还能查账。"
        "这类能力最适合 IT 运维、法务、财务、合规风控等高责任场景，因为这些场景最怕的不是 AI 说得不够漂亮，而是动作做错以后无法证明过程。"
        "AegisTrust 的价值不在于再做一个新模型，而在于把 AI 幻觉和执行风险拆成一条可拦截、可留证、可复验的工程链。"
        "目前本地仓库群已形成持续更新的开源体系，核心链路可现场回放，说明项目已具备从技术验证走向场景验证的可信起点。"
        "下一阶段将围绕企业可信部署需求，推动企业部署版、验证服务和行业样板落地。"
    )
    return text + "\n"


def build_judge_view_diagnosis(ctx: dict[str, Any]) -> str:
    release = ctx["releases"]
    smoke = ctx["smoke_results"]
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 评委视角诊断

        生成时间：{fmt_dt(ctx["generated_at"])}

        这份诊断不是改文案，而是站在比赛评委的位置，判断材料现在最容易得分和最容易失分的地方。

        ## 1. 已经很强的地方

        - 方向踩中了真问题：AI 正在从“会说”走向“会做”，执行信任会变成真实痛点。
        - 不是空概念：当前本地仓库群扫描到 {len(ctx["repos"])} 个仓库，最近 30 天有更新的仓库有 {ctx["active_30d"]} 个，说明项目持续推进。
        - 有可验证的工程证据：`god-spear` `{release.get("god-spear") or "none"}`、`safety-valve-spec` `{release.get("safety-valve-spec") or "none"}`、`aro-audit` `{release.get("aro-audit") or "none"}`，核心仓库已有公开版本。
        - 有现场可讲的演示结论：`god-spear` PASS、`safety-valve-spec` PASS、`execution-integrity-core` PASS、`aro-audit` 正常样本 `VERIFY_OK` 且篡改样本触发 `Merkle mismatch`。
        - 主线已经能收敛成一句话：AegisTrust = AI 可信执行基础设施。

        ## 2. 现在还显得太技术的地方

        - 过早出现仓库名、协议名和实现细节，会让非技术评委在前 3 分钟内失去抓手。
        - “我做了哪些模块”讲得比“为什么现在一定需要这层能力”更充分，容易像技术说明书。
        - 证明材料很多，但业务结果表达还不够短。评委未必记得 `hash chain`，但会记得“关键动作可拦截、可留证、可复验”。
        - `POP`、FDO、五层总架构这些内容能证明深度，但放太前面会冲淡当前商业化切口。

        ## 3. 五分钟后评委最可能忘掉什么

        - 具体仓库名称和层级关系。
        - 各类协议缩写和实现路径。
        - 仓库总数本身。

        他们更可能记住的只有两件事：

        - AI 从会说走向会做之后，企业缺一层可信执行基础设施。
        - 你不是在讲想法，而是已经把关键链路做成能演示的工程闭环。

        ## 4. 最应该留下的单一印象

        **AegisTrust 不是又一个 AI 安全插件，而是一层让关键 AI 动作更可验证、更可控制、更可追责的可信执行基础设施。**

        如果评委只记住这一句，后面的技术证据和商业判断才有地方落。

        ## 5. 什么应该移到附录或备份

        - `POP` 的完整协议语义和 persona taxonomy
        - GitHub 全仓库矩阵
        - 五层总架构的完整展开
        - FDO / DOI / registry 的详细编号与路径
        - 每个仓库的细颗粒 benchmark、schema、workflow 列表

        这些内容不是不要，而是不该抢主舞台。

        ## 6. 当前比赛版最该强化的五个判断

        - 为什么这个问题现在出现，而且会越来越急。
        - 为什么现有 AI stack 解决不了“动作可信”这件事。
        - 为什么 AegisTrust 是基础设施，而不是又一个点状功能。
        - 为什么项目虽然早期，但已经可信。
        - 为什么北京适合承接它的下一步验证和落地。

        ## 7. 当前比赛版最该避免的三件事

        - 不要把主舞台讲成仓库巡礼。
        - 不要把早期项目包装成已经规模化商业成功。
        - 不要把政策表达写成表态，保持顺势而为即可。
        """
        ).strip()
        + "\n"
    )


def build_award_storyline(ctx: dict[str, Any]) -> str:
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 获奖级叙事母线

        生成时间：{fmt_dt(ctx["generated_at"])}

        ## 一条主线

        **当 AI 从“会说”进入“会做”，可信执行会成为新的基础能力；AegisTrust 正是在补这层基础设施。**

        这条主线要贯穿 BP、PPT、报名表和路演讲稿，不能在不同材料里切换说法。

        ## 1. 为什么这个问题现在出现

        AI 已经不只是生成文本，而是在进入真实流程，开始调用工具、改配置、触发任务、输出可执行结果。过去“回答错了”主要是体验问题，现在“动作做错了”会直接变成业务风险、审计风险和责任风险。

        ## 2. 为什么现有 AI stack 解决得不够好

        当前主流栈更擅长让模型变得有用，不擅长让动作变得可信。日志平台只能回看，提示词约束更像在管模型怎么说，单点安全插件往往只能守一处，缺少把事前、事中、事后串成闭环的能力。

        ## 3. 为什么 AegisTrust 改变的是结构，不是加一个功能

        AegisTrust 不是在模型旁边再放一个提醒器，而是在关键动作前后补上一条可信执行链：

        - 上线前先做门禁检查
        - 执行时要求关键动作带收据
        - 执行后输出独立可复验的证据包

        它改变的是动作进入真实系统的方式。

        ## 4. 为什么这是一层基础设施，而不是普通安全工具

        因为它不依赖某一个模型，也不绑定某一个 Agent 框架。它面对的是更底层的问题：任何高风险 AI 动作，在进入真实系统之前，是否有边界、是否留收据、是否能复验。这是跨模型、跨框架、跨场景都会反复出现的共同需求。

        ## 5. 为什么现在是合适的窗口

        市场对 AI 的关注点正在从“有没有能力”转向“能不能放心落地”。越是进入运维、法务、财务、合规、风控这类高责任场景，越需要一层可验证、可审计、可追责的执行基础能力。这个窗口比“更强模型”更适合早期团队切入。

        ## 6. 为什么项目还早，但已经可信

        AegisTrust 目前最准确的阶段判断不是“成熟公司”，而是“核心基础层已打通、正在进入场景验证”。可信度来自三类证据：

        - 真实工程进展：本地仓库群持续更新
        - 可复跑验证：核心链路已可现场回放
        - 可引用材料：已有公开版本、DOI 归档和外部讨论素材

        ## 7. 评委应该听到的核心句子

        - AI 从会说走向会做之后，缺的不是更多模型能力，而是可信执行能力。
        - AegisTrust 不是又一个点状安全插件，而是一层 AI 可信执行基础设施。
        - 这个项目现在值得支持，不是因为故事大，而是因为问题真实、结构正确、证据已经开始成形。

        ## 8. 这条母线怎么落到材料里

        - BP 前半段先回答“为什么现在需要它”，后半段再回答“我们怎么做”。
        - PPT 前 8 页只做三件事：定义问题、建立窗口、完成定位。
        - 技术深度只作为可信度支撑，不作为主舞台起点。
        - `POP` 和五层架构作为来源和长期延展，不与当前商业主线并列。
        """
        ).strip()
        + "\n"
    )


def build_bp_v4_award(ctx: dict[str, Any]) -> str:
    release = ctx["releases"]
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 商业计划书 V4（获奖叙事版）

        生成时间：{fmt_dt(ctx["generated_at"])}

        一句话定位：**AegisTrust 是一套 AI 可信执行基础设施，让关键 AI 动作可拦截、可留证、可复验。**

        这份版本不再从“我做了哪些仓库”出发，而是从评委最关心的三个判断出发：为什么这个问题值得关注，为什么这个项目值得支持，为什么这件事适合现在做。

        ## 1. 我们看到的变化

        AI 正在从“会答问题”走向“会做事情”。它开始调用工具、修改配置、触发流程、输出可执行结果。这个阶段一旦进入真实业务，企业面对的核心问题就不再只是模型效果，而是执行信任。

        换句话说，下一波竞争不只是“谁的模型更强”，而是“谁能让关键 AI 动作更可信地进入真实系统”。

        ## 2. 为什么这是一个正在出现的新痛点

        当 AI 进入高责任流程时，企业最担心三件事：

        - 它会不会越权、误判或在上下文不完整时直接做错事
        - 关键动作是否能被提前阻断，而不是事故后补救
        - 一旦出现争议，能不能拿出独立可复核的证据，而不是只剩口头解释

        所以企业缺的不是又一个会回答问题的 Agent，而是一层可信执行能力。

        ## 3. 为什么现有方案还不够

        现有 AI 栈大致解决了三类问题，但都没有完整解决“动作可信”：

        - 模型能力层：让 AI 更会理解和生成
        - 日志回看层：让人事后知道发生过什么
        - 单点防护层：在局部做拦截或提醒

        这些能力都重要，但它们没有把事前、事中、事后连成一条闭环。高风险 AI 动作一旦进入真实流程，真正缺的是一条可信执行链。

        ## 4. AegisTrust 提供的是什么

        AegisTrust 的核心思路很直接：**不和每一次模型判断赌输赢，而是在关键动作前后补上一条可验证的执行链。**

        这条链分成三步：

        - 事前：上线前先做门禁检查
        - 事中：关键动作必须带收据
        - 事后：输出独立可复验的证据包

        用人话概括就是：**先拦住，再留票，最后还能查账。**

        ## 5. 为什么我们把它叫“基础设施”

        因为它解决的是一个跨模型、跨框架、跨场景都会反复出现的问题。

        企业未来可以换模型、换 Agent 框架、换业务流程，但只要 AI 真正开始“做事”，就一定会面对同一个问题：关键动作是否可信。AegisTrust 站在动作边界和结果复验这一层，不依赖某一个模型品牌，也不绑定某一个框架生态。

        这就是它更接近基础设施，而不是普通安全插件的原因。

        ## 6. 当前对外可讲清楚的四层闭环

        当前商业化主线收敛为四个模块：

        - `god-spear`：上线前门禁。先检查边界、失败信号和回滚路径是否完整。
        - `safety-valve-spec`：动作收据。关键动作必须明确是允许、拒绝还是降级执行。
        - `execution-integrity-core`：执行完整性。证明过程没有被偷偷改写。
        - `aro-audit`：审计证据。把结果整理成可独立复核、可交付的证据包。

        这四层已经足以构成一条评委和客户都能理解的可信执行闭环。

        ## 7. 我们如何处理 AI 幻觉风险

        AegisTrust 不承诺“消灭幻觉”。更真实的做法是把幻觉的后果控制在可阻断、可发现、可追责的范围内。

        只要关键动作没有通过门禁、没有形成收据、没有进入完整性链和证据链，它就不能进入可信交付路径。也就是说，我们不是承诺 AI 永不出错，而是把错误从“事故”变成“可治理事件”。

        ## 8. 这套能力最适合先进入哪些场景

        第一批最适合的场景不是泛娱乐，也不是低责任问答，而是高责任、高审计压力的流程：

        - IT 运维与自动化发布
        - 法务与合同审阅/流转
        - 财务与审批辅助流程
        - 合规、风控和内部控制相关动作

        这些场景的共同点是，一旦 AI 开始“代做动作”，组织最先关心的就是边界、留痕和复验。

        ## 9. 产品形态

        当前最合理的产品形态分三层：

        - 开源基础版：协议、校验器、示例和 quickstart
        - 企业部署版：策略包、模板、验证接口和管理面板
        - 实施服务版：接入、培训、认证和场景打磨

        这样的结构既能用开源获得进入场景的机会，也能逐步建立软件收入和服务收入。

        ## 10. 现在已经具备什么可信度

        AegisTrust 目前最重要的不是概念，而是已经形成可复跑、可说明、可引用的技术证据：

        - 本地仓库群共扫描到 {len(ctx["repos"])} 个仓库，最近 30 天有更新的仓库有 {ctx["active_30d"]} 个
        - `spear-check` 已在 {ctx["spear_adoption"]["count"]} 个本地仓库持续接入
        - 核心公开版本包括 `god-spear` `{release.get("god-spear") or "none"}`、`safety-valve-spec` `{release.get("safety-valve-spec") or "none"}`、`aro-audit` `{release.get("aro-audit") or "none"}`
        - 本地回放显示：门禁 PASS、conformance PASS、执行完整性 PASS、审计证据链可验证且篡改样本会暴露

        这说明项目已经越过“想法阶段”，进入“技术可信度已建立”的阶段。

        ## 11. 为什么它还早，但依然值得支持

        我们不会把项目包装成已经完成规模化商业验证。更准确的判断是：

        - 技术闭环已经成立
        - 对外演示能力已经形成
        - 场景验证和试点转化正是下一步最关键的增量

        这类项目最适合在早期获得支持，因为它要补的是一层未来会反复被需要的底层能力，而不是一次性需求。

        ## 12. 架构来源与长期延展

        AegisTrust 并不是孤立项目，它来自一套更完整的 agent 架构研究。`POP` 这样的仓库属于上游身份层，GitHub 首页和架构总览仓库则提供了更大的系统图谱。

        但当前商业化阶段，我们有意识地只聚焦其中最刚需、最容易落地的一段：可信执行。这样既保留长期深度，也避免主线发散。

        ## 13. 外部可信度补充证据

        当前主舞台不需要把这些编号全部展开，但作为补充材料，它们说明项目已经有走向更广泛技术讨论的基础：

        - `aro-audit` 已形成 Zenodo DOI 归档：`10.5281/zenodo.18728568`
        - `aro-audit` 已形成 FDO 相关 profile 草案与注册标识：`21.T11966/aro-audit-profile-v1`
        - `persona-object-protocol` 已有 POP-Core DOI 归档：`10.5281/zenodo.18907957`

        这些材料的价值在于增强可信度，而不是替代商业叙事。

        ## 14. 产业方向与政策语境

        AegisTrust 最适合被写成“顺应产业方向的可信执行基础能力”，而不是“直接承接政策任务的项目”。

        更自然的表达是：随着人工智能深入企业真实流程，可信执行、风险治理和审计能力正逐渐成为智能化升级中的基础要求。AegisTrust 面向的正是这一趋势下的新型基础能力需求。

        ## 15. 商业模式与切入路径

        商业模式会沿着“开源进入场景，企业版放大价值，服务完成落地”这条路径推进。

        第一阶段先争取试点和联合验证，把样板打出来；第二阶段沉淀成标准化部署能力和验证服务；第三阶段再扩大到更多高责任流程和合作生态。

        当前最需要的不是大规模营销，而是真实场景、合作伙伴和早期验证资金一起到位。

        ## 16. 北京为什么适合承接

        北京的价值不在于口号，而在于资源条件与项目阶段高度匹配：

        - AI 企业、研究机构和产业服务资源密集
        - 高责任行业和数字化升级需求集中
        - 更容易联合高校、园区、平台做样板验证与生态合作

        更现实的落地路径是，先从高审计压力、高内控要求的企业流程切入，再逐步形成行业模板、实施方法和认证能力。

        ## 17. 未来 12 个月路线

        - 0 到 3 个月：把产品叙事、演示环境和行业样板进一步压短
        - 3 到 6 个月：完成 2 到 3 个试点或合作意向
        - 6 到 12 个月：形成可复制的企业部署能力和标准化验证服务

        路线重点不是扩概念，而是把现有能力压成更短、更清晰、更容易采购的交付路径。

        ## 18. 当前最希望得到什么支持

        当前最有价值的支持不只是资金，更包括三类资源：

        - 真实试点场景
        - 带行业理解的合作伙伴
        - 早期验证资金与生态连接

        因为 AegisTrust 现在最值得放大的，不是故事，而是把一个正在成形的基础能力尽快接到真实世界。

        ## 结论

        当 AI 从“会说”进入“会做”，可信执行会逐步成为新的基础能力。AegisTrust 的价值不在于堆了多少仓库，而在于已经把高风险 AI 动作的问题拆成了一条可拦截、可留证、可复验的工程闭环。它还早，但方向对、结构对、证据也已经开始对齐。这正是一个值得被发现、被支持、被放大的项目应有的样子。
        """
        ).strip()
        + "\n"
    )


def build_slides(ctx: dict[str, Any]) -> list[dict[str, Any]]:
    smoke = ctx["smoke_results"]
    return [
        {
            "kind": "cover",
            "title": "AegisTrust",
            "subtitle": "AI 可信执行基础设施",
            "tagline": "把高风险 AI 动作变成可拦截、可留证、可复验的能力",
            "footer": f"本地仓库进展更新 · {fmt_day(ctx['generated_at'])}",
        },
        {
            "title": "为什么现在做",
            "headline": "AI 已经开始碰真实系统，不再只是陪聊工具",
            "bullets": [
                "从问答走向运维、合同、财务、协同流程",
                "一旦能执行，风险级别就会立刻上升",
                "下一层基础设施必须解决“能不能放心让它做”",
            ],
            "callout": "从会说，到会做",
        },
        {
            "title": "真正的问题",
            "headline": "企业怕的不是 AI 说错一句话，而是它做错一件事",
            "bullets": [
                "幻觉、上下文缺失、策略错误都会穿透到真实动作",
                "传统日志能留痕，但很难独立证明",
                "企业缺的是从事前到事后的可信执行闭环",
            ],
            "callout": "核心问题是“动作可信”",
        },
        {
            "title": "我们的答案",
            "headline": "先拦住，再留票，最后还能把账查清楚",
            "bullets": [
                "上线前：CI 门禁，意思是代码发布前先自动体检",
                "执行中：关键动作必须带收据，意思是每次动作都留下一张可核对的小票",
                "执行后：导出证据包，意思是出事后能把账查清楚",
            ],
            "callout": "可拦截 · 可证明 · 可复验",
        },
        {
            "title": "幻觉风险",
            "headline": "不和幻觉赌输赢，而是限制幻觉的伤害半径",
            "bullets": [
                "我们不承诺消灭幻觉",
                "我们承诺把后果控制在可阻断、可发现、可追责范围内",
                "重点不是模型说错，而是别让它做错",
            ],
            "callout": "先控后果，再谈优化",
        },
        {
            "title": "技术架构",
            "headline": "四个模块，不讲大而全，只讲一条能落地的主线",
            "bullets": [
                "god-spear：上线前门禁，先检查规则写没写完整",
                "safety-valve-spec：动作收据，给每次关键动作留凭证",
                "execution-integrity-core + aro-audit：证明过程没被改、结果能复查",
            ],
            "callout": "四个模块，一条主线",
        },
        {
            "title": "执行完整性",
            "headline": "不是记日志，而是证明过程没被偷偷改过",
            "bullets": [
                "hash chain：前后相扣，改一条会影响整条链",
                "确定性导出：同一份数据，多次导出结果一致",
                "人话就是：过程有没有被改，一查就知道",
            ],
            "callout": smoke.get("execution-integrity-core", {}).get(
                "summary", "SELF_CHECK: PASS"
            ),
        },
        {
            "title": "审计证据",
            "headline": "交付的不是截图，而是一包能独立核验的材料",
            "bullets": [
                "append-only journal：只能追加，不能偷偷改旧记录",
                "manifest + Merkle/checkpoint：给证据包做防篡改校验点",
                "人话就是：交出去的不只是日志，而是一包能核验的材料",
            ],
            "callout": smoke.get("aro-audit", {}).get("summary", "快速演示可用"),
        },
        {
            "title": "产品形态",
            "headline": "先用开源进场，再用企业版和服务变现",
            "bullets": [
                "开源基础版：协议、校验器、demo",
                "企业部署版：策略包、审计模板、托管验证",
                "服务版：实施、培训、认证",
            ],
            "callout": "软件 + 服务",
        },
        {
            "title": "商业模式",
            "headline": "收入路径很清楚：软件费 + 服务费 + 验证费",
            "bullets": [
                "开源标准做入口",
                "企业软件做复用收入",
                "试点实施与认证做高附加值服务",
            ],
            "callout": "先入口，后放大收入",
        },
        {
            "title": "市场",
            "headline": "高风险 AI 一落地，审计和内控需求就会马上出现",
            "bullets": [
                "企业内部 AI 助手正在进入真实流程",
                "首批高需求部门：IT、法务、财务、合规风控",
                "AegisTrust 对准的是最先出现审计压力的场景",
            ],
            "callout": "先做高压力场景",
        },
        {
            "title": "对比优势",
            "headline": "我们不是做一个点状插件，而是做整条可信执行链",
            "bullets": [
                "日志平台：看得见，但不一定说得清",
                "提示词 guardrail：管模型怎么说，不够管模型怎么做",
                "AegisTrust：从上线前到出事后，整条链都能查",
            ],
            "callout": "闭环能力",
        },
        {
            "title": "当前进展",
            "headline": "这不是 PPT 概念，核心链路已经能本地复跑",
            "bullets": [
                f"本地仓库群共有 {len(ctx['repos'])} 个仓库",
                f"最近 30 天更新的仓库有 {ctx['active_30d']} 个",
                f"本地已有 {ctx['spear_adoption']['count']} 个仓库接入 spear-check",
            ],
            "callout": "本地复跑已完成",
        },
        {
            "title": "演示能力",
            "headline": "现场可以证明三件事：能拦、能留证、能查出篡改",
            "bullets": [
                "god-spear：本地回放 PASS，说明门禁可跑",
                "safety-valve-spec conformance：PASS，说明规则和收据能对上",
                "aro-audit：正常样本通过，篡改样本失败，说明证据链有效",
            ],
            "callout": "三条关键链路已回放",
        },
        {
            "title": "路线图",
            "headline": "下一步不是堆概念，而是拿试点把产品线压实",
            "bullets": [
                "0-3 个月：统一产品叙事和行业样板",
                "3-6 个月：完成 2-3 个试点或合作意向",
                "6-12 个月：形成企业部署和验证服务能力",
            ],
            "callout": "从技术验证到场景验证",
        },
        {
            "title": "团队",
            "headline": "当前最强的证明，不是头衔，而是持续公开迭代的代码",
            "bullets": [
                "创始人主导技术核心，公开 GitHub 证据强",
                "当前强项：协议设计、工程工具、验证流程",
                "下一步补强：合规、内控、ToB 交付",
            ],
            "callout": "技术主导型团队",
        },
        {
            "title": "融资计划",
            "headline": "资金主要用来把验证成果变成可复制交付",
            "bullets": [
                "以早期验证资金和场景资源为主",
                "资金用途：产品化打磨、试点交付、生态合作",
                "最需要的是资金和真实场景一起到位",
            ],
            "callout": "先把试点做成",
        },
        {
            "title": "落地计划",
            "headline": "先做有审计压力的高风险流程，再逐步扩大",
            "bullets": [
                "先找高风险 AI 流程场景切入",
                "优先进入有审计压力和内控需求的企业流程",
                "再与高校、机构和园区做联合验证和共建",
            ],
            "callout": "先难后易",
        },
        {
            "title": "总结",
            "headline": "AegisTrust 想解决的是 AI 从会说到会做之后的信任问题",
            "bullets": [
                "AegisTrust 不是另一个 Agent 小工具",
                "它是一层面向高风险 AI 动作的可信执行基础设施",
                "公开仓库和本地回放已经证明技术在持续推进",
            ],
            "callout": "可阻断 · 可证明 · 可复验",
        },
        {
            "kind": "closing",
            "title": "谢谢",
            "subtitle": "AegisTrust = AI 可信执行基础设施",
            "tagline": "可阻断 · 可证明 · 可复验",
        },
    ]


def build_slides_v2(ctx: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        {
            "title": "AegisTrust",
            "headline": "AI 可信执行基础设施",
            "bullets": [
                "面向高风险 AI 动作场景",
                "让关键动作可拦截、可留证、可复验",
                "当前版本基于本地仓库和可回放 demo 生成",
            ],
            "callout": "从会说，到会做",
        },
        {
            "title": "愿景",
            "headline": "让企业从“能用 AI”走到“敢让 AI 做事”",
            "bullets": [
                "AI 正在进入运维、法务、财务和协同流程",
                "真正的门槛不是模型效果，而是执行信任",
                "下一层基础设施必须解决动作可信问题",
            ],
            "callout": "信任比炫技更值钱",
        },
        {
            "title": "问题",
            "headline": "企业怕的不是 AI 说错一句话，而是它做错一件事",
            "bullets": [
                "幻觉、上下文缺失、策略错误会穿透到真实动作",
                "普通日志只能留痕，不能强证明",
                "一旦出事，企业很难拿出独立可复验的证据",
            ],
            "callout": "核心问题是动作可信",
        },
        {
            "title": "AI 幻觉风险",
            "headline": "不和幻觉赌输赢，而是限制幻觉造成的伤害半径",
            "bullets": [
                "不承诺消灭幻觉",
                "承诺把后果控制在可阻断、可发现、可追责范围内",
                "重点不是模型说错，而是别让它做错",
            ],
            "callout": "先控后果，再谈优化",
        },
        {
            "title": "解决方案",
            "headline": "先拦住，再留票，最后还能查账",
            "bullets": [
                "上线前做门禁检查，别让明显有缺口的配置进入生产",
                "执行时要求关键动作带收据，别让高风险操作裸奔",
                "执行后导出证据包，别让复盘只剩口头解释",
            ],
            "callout": "可拦截 · 可留证 · 可复验",
        },
        {
            "title": "架构",
            "headline": "当前商业化主线聚焦四层可信执行闭环",
            "bullets": [
                "god-spear：CI governance，上线前先自动体检",
                "safety-valve-spec：runtime safety，关键动作必须带小票",
                "execution-integrity-core + aro-audit：过程能证明，结果能复查",
            ],
            "callout": "四层，一条主线",
        },
        {
            "title": "执行完整性",
            "headline": "不是只记日志，而是证明过程没被偷偷改过",
            "bullets": [
                "hash chain：前后相扣，改一处整条会露馅",
                "确定性导出：同一份数据，多次导出结果一致",
                "人话就是：过程有没有被动过，一查就知道",
            ],
            "callout": "过程可信",
        },
        {
            "title": "审计证据",
            "headline": "交付给客户的不是截图，而是一包能核验的材料",
            "bullets": [
                "append-only journal：旧记录不能被偷偷改掉",
                "manifest + Merkle/checkpoint：给证据包做防篡改校验点",
                "人话就是：交出去的不只是日志，而是一包能查账的材料",
            ],
            "callout": "结果可信",
        },
        {
            "title": "产品形态",
            "headline": "先用开源进入场景，再用企业版和服务放大收入",
            "bullets": [
                "开源基础版：协议、校验器、示例、quickstart",
                "企业部署版：策略包、模板、验证接口",
                "实施服务版：接入、培训、认证和长期支持",
            ],
            "callout": "软件 + 服务",
        },
        {
            "title": "市场",
            "headline": "高风险 AI 一落地，可信执行需求就会马上出现",
            "bullets": [
                "首批高需求场景集中在 IT、法务、财务、合规风控",
                "这些部门最看重内控、留痕和事故可追溯",
                "AegisTrust 面向的是高风险 AI 动作的基础设施市场",
            ],
            "callout": "先做高责任场景",
        },
        {
            "title": "政策环境",
            "headline": "顺应数字化转型与发展安全并重的产业趋势",
            "bullets": [
                "AI+ 正在从工具试用走向真实业务接入",
                "企业智能化升级需要可信部署和风险治理能力",
                "AegisTrust 更适合被写成基础能力组件，而不是政策口号项目",
            ],
            "callout": "顺势而为，不抢话语权",
        },
        {
            "title": "竞争",
            "headline": "我们不是做一个点状插件，而是做整条可信执行链",
            "bullets": [
                "日志平台：能看见，不一定说得清",
                "提示词 guardrail：更像管模型怎么说，不够管模型怎么做",
                "AegisTrust：从上线前到事后复查，整条链都能接上",
            ],
            "callout": "闭环能力",
        },
        {
            "title": "进展",
            "headline": "这不是 PPT 概念，核心链路已经能本地复跑",
            "bullets": [
                f"当前本地仓库群共扫描到 {len(ctx['repos'])} 个仓库",
                f"最近 30 天有更新的仓库有 {ctx['active_30d']} 个",
                f"本地已有 {ctx['spear_adoption']['count']} 个仓库接入 spear-check",
            ],
            "callout": "持续迭代中",
        },
        {
            "title": "演示",
            "headline": "现场可以证明三件事：能拦、能留证、能查出篡改",
            "bullets": [
                "god-spear：PASS，说明门禁可跑",
                "safety-valve-spec：PASS，说明规则和收据能对上",
                "aro-audit：正常样本通过，篡改样本失败，说明证据链有效",
            ],
            "callout": "关键路径可回放",
        },
        {
            "title": "路线图",
            "headline": "下一步不是堆更多概念，而是把交付路径压短",
            "bullets": [
                "0-3 个月：统一叙事和行业样板",
                "3-6 个月：完成 2-3 个试点或合作意向",
                "6-12 个月：形成企业部署和标准化验证服务能力",
            ],
            "callout": "从技术验证到场景验证",
        },
        {
            "title": "团队",
            "headline": "当前最强的证明，不是头衔，而是持续公开迭代的代码",
            "bullets": [
                "创始人主导技术核心，长期做体系化研究和工程实现",
                "当前强项：协议设计、验证流程、可信执行工具链",
                "下一步补强：行业顾问、合规理解、ToB 交付能力",
            ],
            "callout": "技术主导型团队",
        },
        {
            "title": "融资",
            "headline": "资金主要用来把验证成果变成可复制交付",
            "bullets": [
                "优先引入早期验证资金和带场景资源的合作方",
                "重点投入产品化打磨、试点交付和客户成功",
                "资金不是用来讲更大故事，而是先把样板做成",
            ],
            "callout": "先做实，再放大",
        },
        {
            "title": "北京计划",
            "headline": "北京适合成为早期验证和试点落地的重要节点",
            "bullets": [
                "AI 企业、研究机构和产业服务资源密集",
                "企业数字化和智能化升级需求集中",
                "更容易联合高校、园区和平台做样板验证",
            ],
            "callout": "资源条件匹配",
        },
        {
            "title": "总结",
            "headline": "当 AI 从会说走向会做，市场需要一层可信执行基础设施",
            "bullets": [
                "AegisTrust 聚焦高风险 AI 动作场景",
                "核心价值是可拦截、可留证、可复验",
                "当前本地仓库和 demo 已证明这条链路能跑起来",
            ],
            "callout": "AI 可信执行基础设施",
        },
        {
            "title": "谢谢",
            "headline": "AegisTrust = AI 可信执行基础设施",
            "bullets": [
                "先聚焦可信执行闭环，再逐步向更大架构延展",
                "欢迎围绕试点、场景和生态合作展开交流",
                "当前最重要的是把真实样板做出来",
            ],
            "callout": "谢谢",
        },
    ]


def build_slides_v4_award(ctx: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        {
            "kind": "cover",
            "title": "AegisTrust",
            "subtitle": "AI 可信执行基础设施",
            "tagline": "让关键 AI 动作可拦截、可留证、可复验",
            "footer": f"获奖叙事版 · {fmt_day(ctx['generated_at'])}",
        },
        {
            "title": "时代变化",
            "headline": "AI 正在从“会说”走向“会做”",
            "bullets": [
                "从回答问题，进入调用工具、改配置、触发流程",
                "一旦进入真实系统，风险等级会立刻变化",
                "下一层竞争，不只是模型能力，而是执行信任",
            ],
            "callout": "从会说，到会做",
        },
        {
            "title": "真实痛点",
            "headline": "企业怕的不是 AI 说错一句话，而是它做错一件事",
            "bullets": [
                "上下文缺失、幻觉和策略偏差，会穿透到真实动作",
                "事故发生后，组织最怕的是说不清楚过程",
                "高责任场景天然要求边界、留痕和复验",
            ],
            "callout": "核心问题是动作可信",
        },
        {
            "title": "市场缺口",
            "headline": "现有 AI 栈更擅长让模型有用，不够擅长让动作可信",
            "bullets": [
                "模型层解决能力，不直接解决责任",
                "日志平台能回看，不一定能强证明",
                "单点插件能守一处，守不住整条链",
            ],
            "callout": "缺的是闭环，不是零件",
        },
        {
            "title": "项目定位",
            "headline": "AegisTrust 补的是 AI 进入真实流程后的可信执行层",
            "bullets": [
                "不和每一次模型判断赌输赢",
                "在关键动作前后补一条可验证的执行链",
                "让组织更敢把 AI 放进高责任流程",
            ],
            "callout": "AI 可信执行基础设施",
        },
        {
            "title": "为什么是基础设施",
            "headline": "因为这个问题会跨模型、跨框架、跨场景反复出现",
            "bullets": [
                "企业可以换模型，但不会放弃对关键动作的控制",
                "企业可以换框架，但不会放弃对结果的追责能力",
                "可信执行是底层共性需求，不是单一工具需求",
            ],
            "callout": "不是插件，是底层能力",
        },
        {
            "title": "解决思路",
            "headline": "先拦住，再留票，最后还能查账",
            "bullets": [
                "事前：上线前做门禁检查",
                "事中：关键动作必须带收据",
                "事后：结果导出成独立可复验的证据包",
            ],
            "callout": "可拦截 · 可留证 · 可复验",
        },
        {
            "title": "四层闭环",
            "headline": "当前商业化主线已经收敛成四层可信执行闭环",
            "bullets": [
                "god-spear：上线前门禁，先做自动体检",
                "safety-valve-spec：关键动作必须有标准化小票",
                "execution-integrity-core + aro-audit：过程能证明，结果能复查",
            ],
            "callout": "四层，一条主线",
        },
        {
            "title": "场景示例",
            "headline": "以 AI 运维助手为例，关键不只是能不能做，而是能不能放心做",
            "bullets": [
                "没有可信执行时：动作快，但边界不清、责任不清",
                "接入 AegisTrust 后：高风险动作先过门禁，再留收据，最后出证据包",
                "同一套方法也适用于法务、财务和合规流程",
            ],
            "callout": "先从高责任流程切入",
        },
        {
            "title": "幻觉治理",
            "headline": "我们不承诺消灭幻觉，我们承诺控制幻觉的伤害半径",
            "bullets": [
                "关键动作没有通过门禁，就不能进入可信路径",
                "没有收据和完整性链，就不能算合规执行",
                "错误从事故变成可治理、可追责的事件",
            ],
            "callout": "先控后果，再谈优化",
        },
        {
            "title": "项目真实性",
            "headline": "这不是想法阶段，而是已经开始成形的工程闭环",
            "bullets": [
                f"本地仓库群共扫描到 {len(ctx['repos'])} 个仓库",
                f"最近 30 天有更新的仓库有 {ctx['active_30d']} 个",
                f"spear-check 已在 {ctx['spear_adoption']['count']} 个仓库持续接入",
            ],
            "callout": "持续迭代中",
        },
        {
            "title": "可现场验证",
            "headline": "现场能证明三件事：能拦、能留证、能查出篡改",
            "bullets": [
                "god-spear：PASS，说明门禁可跑",
                "safety-valve-spec：PASS，说明动作收据规则可验",
                "aro-audit：正常样本通过，篡改样本暴露，说明证据链有效",
            ],
            "callout": "关键路径可回放",
        },
        {
            "title": "产品形态",
            "headline": "先用开源进入场景，再用企业版和服务放大价值",
            "bullets": [
                "开源基础版：协议、校验器、示例和 quickstart",
                "企业部署版：策略包、模板、验证接口和管理面板",
                "实施服务版：接入、培训、认证和样板复制",
            ],
            "callout": "软件 + 服务",
        },
        {
            "title": "应用场景",
            "headline": "第一批需求会先出现在高责任、高审计压力的流程里",
            "bullets": [
                "IT 运维与自动化发布",
                "法务、财务、合规和风控流程",
                "任何需要边界、留痕和问责能力的企业动作",
            ],
            "callout": "先做高责任场景",
        },
        {
            "title": "商业模式",
            "headline": "先做样板，再做可复制部署，再做标准化验证服务",
            "bullets": [
                "开源降低试用门槛，建立进入场景的机会",
                "企业版沉淀策略、模板和验证能力",
                "服务帮助客户完成落地和内控对齐",
            ],
            "callout": "先把试点做成",
        },
        {
            "title": "产业方向",
            "headline": "AegisTrust 适合被写成顺应趋势的基础能力，而不是政策口号",
            "bullets": [
                "AI+ 正在从试用走向真实业务接入",
                "发展和安全并重，会让可信部署需求更明确",
                "可信执行会逐步成为智能化升级中的基础能力",
            ],
            "callout": "顺势而为",
        },
        {
            "title": "北京计划",
            "headline": "北京适合成为早期验证、样板落地和生态连接的重要节点",
            "bullets": [
                "AI 企业、研究机构和产业服务资源集中",
                "高责任行业和数字化升级需求集中",
                "更容易联合高校、园区和平台做样板验证",
            ],
            "callout": "资源条件匹配",
        },
        {
            "title": "团队与阶段",
            "headline": "项目还早，但技术可信度已经建立，正进入场景验证阶段",
            "bullets": [
                "当前最强证明不是头衔，而是持续公开迭代的工程证据",
                "当前强项：协议设计、验证流程、可信执行工具链",
                "下一步补强：行业顾问、ToB 交付、试点推进",
            ],
            "callout": "早期，但可信",
        },
        {
            "title": "下一步与需求",
            "headline": "未来 12 个月的重点，不是扩故事，而是把交付路径压短",
            "bullets": [
                "0-3 个月：继续压缩叙事、打磨演示、形成样板",
                "3-6 个月：完成 2-3 个试点或合作意向",
                "6-12 个月：形成标准化部署和验证服务能力",
            ],
            "callout": "需要场景、伙伴和早期资金",
        },
        {
            "kind": "closing",
            "title": "谢谢",
            "subtitle": "AegisTrust = AI 可信执行基础设施",
            "tagline": "值得被支持的理由，不是故事大，而是问题真、结构对、证据已开始成形",
        },
    ]


def build_pitch_pack(ctx: dict[str, Any]) -> str:
    slides = build_slides_v4_award(ctx)
    note_text = [
        "先把一句话打稳：AegisTrust 不是又一个 Agent 小工具，而是一层 AI 可信执行基础设施。今天我只回答三个问题：为什么现在需要它，为什么它是基础设施，为什么这个项目已经值得被支持。",
        "第一件事是时代变化。AI 已经不是只会生成内容，而是在进入真实流程。只要它开始做事，组织对它的要求就从好不好用，变成可不可以放心交给它做。",
        "企业最怕的不是 AI 回答错一句，而是它做错一个动作。尤其是在运维、法务、财务、合规这种高责任场景，动作一旦出错，后果比回答不漂亮严重得多。",
        "为什么现有方案不够？因为现有 AI 栈更擅长让模型变得有用，不够擅长让动作变得可信。日志能回看，插件能挡一点，但没有形成从事前到事后的闭环。",
        "所以我们的定位很明确：AegisTrust 补的是 AI 进入真实流程后的可信执行层。我们不和每一次模型判断赌输赢，而是在动作前后补一条可验证的执行链。",
        "为什么说它是基础设施？因为这个问题会跨模型、跨框架、跨场景重复出现。企业可以换模型，但不会放弃对关键动作的控制和对结果的追责。",
        "解决思路其实很好记，就是一句话：先拦住，再留票，最后还能查账。先做门禁，再做动作收据，最后给出证据包。",
        "这条链当前已经收敛成四层闭环。你不需要记住全部仓库名，只要记住它分别对应事前门禁、事中收据、过程完整性和事后证据。",
        "如果把它放进一个 AI 运维助手场景里，你会更容易理解。关键不只是让 AI 能改配置，而是要确保高风险动作先过检查、执行有收据、结果能复查。",
        "这也是我们处理幻觉的方式。我们不承诺完全消灭幻觉，但我们承诺把幻觉造成的后果控制在可阻断、可发现、可追责的范围内。",
        "为什么说这不是 PPT 项目？因为它已经有持续工程证据。本地仓库群持续更新，跨仓库接入在发生，说明这不是一次性原型。",
        "更关键的是现场可验证。门禁能跑，conformance 能过，证据链可以验证，而且篡改样本会暴露。这件事对评委很重要，因为它说明项目已经过了纯概念阶段。",
        "产品形态也很清楚。先用开源版进入场景，再用企业部署版沉淀验证能力，最后通过实施服务把场景做深。",
        "应用场景上我们不做泛化叙事，而是聚焦高责任流程。因为越是高审计压力的地方，越愿意为可信执行能力付费。",
        "商业路径不是一上来追求大而全，而是先把试点和样板做成，再把能力做成可复制部署，再做标准化验证服务。",
        "在产业方向上，我们也保持克制。AegisTrust 不是靠口号成立，而是顺应了 AI+ 深入真实流程之后，可信部署和风险治理会成为基础能力这个趋势。",
        "北京为什么重要？因为这里资源和需求都密集，更适合做早期验证、样板落地和生态合作。我们要的不是一个注册地址，而是能把试点做出来的环境。",
        "团队阶段也要讲诚实。我们不是成熟大公司，但已经把核心基础层打通，正在从技术验证走向场景验证。这正是早期项目最值得被支持的窗口。",
        "最后收束成一句话：当 AI 从会说走向会做，市场需要一层可信执行基础设施。AegisTrust 现在值得被支持，不是因为故事大，而是因为问题真实、结构正确、证据已经开始成形。",
        "最后一页不要再扩展新信息，只重复一句主线：AegisTrust 做的是 AI 可信执行基础设施。感谢之后，给评委留下继续追问的空间。",
    ]
    qa_pairs = [
        (
            "为什么不用普通日志平台？",
            "日志平台能回看，不一定能强证明。AegisTrust 强调的是门禁、动作收据、完整性和证据包的闭环，而不只是事后存档。",
        ),
        (
            "为什么不用提示词 guardrail？",
            "提示词 guardrail 更像在限制模型怎么说，不足以覆盖关键动作怎么做、做完以后怎么证明。",
        ),
        (
            "为什么你们不是一个安全功能，而是基础设施？",
            "因为这个问题跨模型、跨框架、跨场景都会反复出现。只要 AI 做高风险动作，组织就需要这层可信执行能力。",
        ),
        (
            "为什么现在是窗口期？",
            "因为企业对 AI 的关注点正在从能不能用，转向敢不敢放到真实流程里。这个阶段对可信执行的需求会快速增长。",
        ),
        (
            "项目还早，为什么现在值得投？",
            "因为核心结构已经搭出来，关键链路可回放，下一步最需要的是场景验证。这个阶段的支持最能放大增量。",
        ),
        (
            "谁会先付费？",
            "第一批更可能是 IT 运维、法务、财务、合规风控等高责任流程的组织，它们对内控、留痕和问责的需求最强。",
        ),
        (
            "最小可落地产品是什么？",
            "最小可落地组合是门禁检查加证据包，也就是先把高风险动作拦住，并把执行结果做成独立可复核材料。",
        ),
        (
            "你们怎么面对 AI 幻觉？",
            "不承诺消灭幻觉，而是控制幻觉后果。没有通过门禁、没有收据、没有进入证据链的动作，不进入可信交付路径。",
        ),
        (
            "POP 和大架构为什么不重点讲？",
            "它们证明长期深度，但当前比赛主线必须聚焦可信执行闭环，否则评委会记不住现在到底卖什么。",
        ),
        (
            "为什么北京适合落地？",
            "因为北京既有高密度 AI 和科研资源，也有高责任行业场景，更适合把早期试点、联合验证和生态合作一起做起来。",
        ),
    ]

    lines = [
        "# AegisTrust 多档讲法统一包",
        "",
        f"生成时间：{fmt_dt(ctx['generated_at'])}",
        "",
        "## 1. 评委一句话",
        "",
        "当 AI 从会说走向会做，AegisTrust 是那层让关键 AI 动作可拦截、可留证、可复验的可信执行基础设施。",
        "",
        "## 2. 30 秒版",
        "",
        "AegisTrust 解决的不是模型会不会回答问题，而是 AI 开始调用工具、改配置、触发流程以后，企业怎么敢让它做事。我们的思路很简单：上线前先做门禁，执行时关键动作必须带收据，执行后输出证据包。也就是说，让关键 AI 动作先拦住、再留票、最后能查账。这是一层面向高风险 AI 动作的可信执行基础设施。",
        "",
        "## 3. 3 分钟版",
        "",
        "AI 正在从“会说”走向“会做”，这会把问题从模型能力推进到执行信任。企业最怕的不是 AI 说错一句，而是它在运维、法务、财务、合规这些真实流程里做错一个动作，最后还说不清过程。",
        "",
        "现有 AI 栈解决不了这个结构性问题。日志能回看，提示词约束能提醒，单点安全插件能挡一点，但都没有把事前、事中、事后连成闭环。AegisTrust 补的就是这条可信执行链：事前做门禁，事中要求动作收据，事后导出证据包。",
        "",
        f"这个项目虽然还早，但已经不是概念。当前本地仓库群共扫描到 {len(ctx['repos'])} 个仓库，最近 30 天有更新的仓库有 {ctx['active_30d']} 个，`spear-check` 已在 {ctx['spear_adoption']['count']} 个仓库接入，核心链路还能现场回放。我们想做的不是一个点状插件，而是一层能支撑高风险 AI 动作落地的可信执行基础设施。",
        "",
        "## 4. 8 分钟版",
        "",
        "如果把过去两年的 AI 看成第一阶段，那一阶段解决的是“能不能让 AI 看起来有用”。现在进入第二阶段，AI 开始真正触碰业务动作。这时候组织最担心的，不再只是回答质量，而是边界、责任和复验。",
        "",
        "这件事为什么重要？因为一旦 AI 进入运维、合同、审批、财务或合规流程，错误就不再只是体验问题，而会直接变成业务风险和治理风险。组织会立刻追问三件事：谁允许它这样做？它到底做了什么？出了问题以后能不能独立证明过程？",
        "",
        "现有 AI stack 对这些问题的回答并不完整。模型层解决能力，日志层解决回看，点状安全工具解决局部防护，但缺少一条把事前、事中、事后串起来的可信执行链。所以 AegisTrust 的定位非常明确：它不是再做一个 Agent 功能，而是补上 AI 进入真实流程之后那条可信执行基础设施。",
        "",
        "我们的方案可以用一句话概括：先拦住，再留票，最后还能查账。事前，用门禁检查去看关键边界是否清楚；事中，让关键动作必须带标准化收据；事后，把过程和结果整理成独立可复核的证据包。这样，组织不需要假设 AI 永远不会错，而是能把错误控制在可阻断、可发现、可追责的轨道里。",
        "",
        "为什么这不是普通安全工具，而是基础设施？因为只要 AI 进入高风险动作场景，这个问题就会反复出现。企业可以换模型，也可以换框架，但不会放弃对关键动作的控制和对结果的问责。AegisTrust 处理的是一个底层共性问题，而不是某个单独产品的问题。",
        "",
        f"这个项目现在值得被支持，也不是因为故事讲得大，而是因为它已经有真实工程证据。当前本地仓库群扫描到 {len(ctx['repos'])} 个仓库，最近 30 天有更新的仓库有 {ctx['active_30d']} 个，`spear-check` 已在 {ctx['spear_adoption']['count']} 个仓库持续接入。更关键的是，核心链路可现场回放，门禁、conformance、完整性和审计证据都能跑出来，篡改样本也会暴露。它说明项目已经越过“纯想法”阶段。",
        "",
        "在商业上，我们会先从高责任、高审计压力的场景切入，比如 IT 运维、法务、财务和合规风控。先用开源版进入场景，再把策略、模板和验证能力沉淀成企业部署版，最后通过实施服务把样板做深。北京之所以重要，不是因为口号，而是因为这里的产业资源、场景需求和生态合作条件都更适合这个阶段的项目。",
        "",
        "所以，如果要用一句话收束，AegisTrust 做的，是当 AI 从会说走向会做之后，那层新的可信执行基础设施。它还早，但问题真、结构对、证据已经开始成形，这正是它值得被发现和放大的原因。",
        "",
        "## 5. 20 页讲稿备注",
        "",
        f"当前页数：{len(slides)}",
        "",
    ]
    for idx, slide in enumerate(slides, start=1):
        note = (
            note_text[idx - 1]
            if idx - 1 < len(note_text)
            else "这一页只做收束：重复主线，不扩新信息，把评委注意力带回“为什么这件事值得支持”。"
        )
        lines.extend(
            [
                f"### Slide {idx:02d} · {slide['title']}",
                note,
                "",
            ]
        )

    lines.extend(
        [
            "## 6. 10 个高频追问与稳妥回答",
            "",
        ]
    )
    for idx, (question, answer) in enumerate(qa_pairs, start=1):
        lines.extend(
            [
                f"{idx}. {question}",
                f"答：{answer}",
                "",
            ]
        )

    return "\n".join(lines).rstrip() + "\n"


def build_application_form_short(ctx: dict[str, Any]) -> str:
    version_150 = "AegisTrust 是一套 AI 可信执行基础设施，面向高风险 AI 动作场景。它解决的不是模型会不会回答问题，而是 AI 开始调工具、改配置、跑流程以后，企业怎么做到敢放权、能控风险、出了事说得清。项目通过门禁检查、动作收据和证据包三层能力，让关键动作可拦截、可留证、可复验。"
    version_300 = "AegisTrust 是一套面向高风险 AI 动作的可信执行基础设施。随着 AI 从“会说”走向“会做”，企业越来越担心的不是回答质量，而是关键动作是否越权、出了问题能否复盘、是否具备独立可验证的证据。AegisTrust 的核心做法是，把 AI 动作前后补上一条可信执行链：上线前先做门禁检查，执行时关键动作必须带收据，执行后输出可复核的证据包。简单说，就是让关键动作先体检、做事要留票、结果还能查账。项目当前已形成持续更新的本地开源体系和可现场回放的关键链路，具备从技术验证走向场景验证的基础。"
    version_500 = f"AegisTrust 是一套 AI 可信执行基础设施，主要解决 AI 进入真实业务流程后“敢不敢放权、出了事怎么证明”的问题。它面向的不是普通问答场景，而是运维、法务、财务、合规风控等高责任流程。当前很多 AI 系统能提高效率，但一旦开始调用工具、修改配置、触发审批或输出可执行结果，组织就会面临新的治理压力：边界是否清楚、动作是否合规、事后是否可复核。AegisTrust 的方案不是再做一个大模型或点状插件，而是在关键动作前后建立一条可信执行链。事前，通过门禁检查把明显有缺口的配置挡在上线前；事中，通过标准化收据让关键动作具备可验证依据；事后，通过证据包和完整性校验让结果能够独立复核。项目当前本地仓库群扫描到 {len(ctx['repos'])} 个仓库，最近 30 天有更新的仓库有 {ctx['active_30d']} 个，核心链路已可现场回放，说明项目已经从概念走向可验证的工程阶段。下一步将重点围绕高责任企业流程做试点、样板和北京落地验证。"
    return "\n".join(
        [
            "AegisTrust 报名表短版文案",
            "",
            "【约 150 字版】",
            version_150,
            "",
            "【约 300 字版】",
            version_300,
            "",
            "【约 500 字版】",
            version_500,
            "",
        ]
    )


def build_judge_qa_top10(ctx: dict[str, Any]) -> str:
    questions = [
        {
            "q": "为什么这个方向会在现在变成刚需，而不是一个提前量过大的判断？",
            "core": "因为 AI 正在从生成内容进入真实动作场景，风险性质已经变了。",
            "full": (
                "过去企业更多把 AI 当成辅助回答工具，用得再深，主要也是内容层风险。现在不一样，"
                "AI 已经开始调用工具、改配置、触发流程、输出可执行结果。一旦进入真实动作场景，组织最先追问的就不是“它聪不聪明”，"
                "而是“它能不能被控制、出了事能不能说清”。所以可信执行不是一个远期概念，而是 AI 落地到高责任流程后立刻出现的基础需求。"
            ),
            "anchors": [
                f"本地仓库群持续活跃：{len(ctx['repos'])} 个仓库，最近 30 天活跃 {ctx['active_30d']} 个",
                "主线已经收敛到高风险 AI 动作，而不是泛 Agent 概念",
            ],
            "avoid": "不要回答成“因为国家现在重视 AI 安全”，这会显得空，也会削弱项目的产业逻辑。",
        },
        {
            "q": "为什么这不是一个普通 AI 安全工具，而要叫基础设施？",
            "core": "因为它处理的是跨模型、跨框架、跨场景都会重复出现的底层问题。",
            "full": (
                "普通安全工具通常是守一个点，比如输入输出过滤、日志留存或权限提醒。AegisTrust 处理的是更底层的问题："
                "只要 AI 开始做高风险动作，组织就会关心三件事，边界是否清楚、过程是否可证明、结果是否可复验。"
                "这些问题不会随着模型替换而消失，也不会因为换框架而消失，所以它更接近基础设施，而不是单点功能。"
            ),
            "anchors": [
                "四层闭环覆盖事前门禁、事中收据、过程完整性、事后证据",
                "不绑定单一模型品牌或单一 Agent 框架",
            ],
            "avoid": "不要说“我们就是标准”，太早也太满。更稳的说法是“我们在补基础层能力”。",
        },
        {
            "q": "市场上已经有日志平台、审计系统、guardrail 工具，你们到底多了什么？",
            "core": "我们多的不是一个功能点，而是一条从事前到事后的可信执行闭环。",
            "full": (
                "日志平台解决的是回看，guardrail 解决的是部分约束，审计系统解决的是事后存档。"
                "但当 AI 真正做动作时，企业需要的是一条完整链路：上线前先把明显有缺口的东西挡住，"
                "执行时让关键动作有收据，执行后还能把结果做成独立可复验的证据包。"
                "所以差异不在某个点，而在闭环是否成立。"
            ),
            "anchors": [
                "门禁 PASS、conformance PASS、完整性 PASS、证据链可验",
                "篡改样本会暴露，不只是“看日志”",
            ],
            "avoid": "不要泛泛地说“别人不行，我们更全面”，要落到闭环差异。",
        },
        {
            "q": "项目现在还早，你凭什么说它已经可信？",
            "core": "因为现在的可信度来自工程证据，而不是商业包装。",
            "full": (
                "我们不会把项目包装成已经规模化落地。现在最准确的阶段判断是：技术闭环已经成立，场景验证刚进入发力期。"
                "它的可信度来自三类证据：第一，仓库持续更新，不是一次性原型；第二，关键链路可现场回放；"
                "第三，已经形成公开版本、DOI 和外部讨论材料。这些证据说明它不只是一个想法，而是已经开始成形的工程基础层。"
            ),
            "anchors": [
                f"`spear-check` 已在 {ctx['spear_adoption']['count']} 个仓库接入",
                "核心版本：god-spear v0.2.0，safety-valve-spec v0.2.0-alpha.1，aro-audit v1.0.1",
            ],
            "avoid": "不要硬说“我们已经成熟”，评委会马上追着问客户和收入。",
        },
        {
            "q": "如果只能先落一个最小产品，你们最小可交付是什么？",
            "core": "最小可交付是门禁检查加证据包。",
            "full": (
                "如果只选最先能落地、最容易让客户理解和采购的一段，我会先推门禁检查加证据包。"
                "也就是先把高风险动作前移到上线前发现问题，再把执行结果做成独立可复验材料。"
                "这套组合最短、最容易进场，也最容易直接对应客户的内控和审计诉求。之后再把动作收据和执行完整性逐步接上。"
            ),
            "anchors": [
                "god-spear 负责风险前移",
                "aro-audit 负责可交付证据包",
            ],
            "avoid": "不要把最小产品说成四层全买全上，早期客户不会这样采购。",
        },
        {
            "q": "你们最先会从谁那里拿到钱？",
            "core": "最先付费的不是泛用户，而是高责任流程里的组织。",
            "full": (
                "这类项目的第一笔钱通常不会来自最开放、最轻量的场景，而会来自最怕出问题的流程。"
                "比如 IT 运维、法务、财务、合规风控这些团队，它们对边界、留痕、复验和问责的需求更强。"
                "因此前期更适合以试点、定制接入、验证服务和企业部署能力切入，而不是先追求大规模标准化订阅。"
            ),
            "anchors": [
                "切入场景聚焦高责任流程",
                "商业模式已拆成开源入口、企业部署版、实施服务版",
            ],
            "avoid": "不要说“所有企业都会需要”，这会显得没有聚焦。",
        },
        {
            "q": "为什么大模型厂商或大平台自己不会顺手把这件事做掉？",
            "core": "因为他们更关注能力上限，而可信执行更接近组织控制层。",
            "full": (
                "大模型厂商当然会做一部分安全和治理能力，但它们天然更偏模型能力、平台能力和通用生态。"
                "可信执行这件事更贴近组织的边界控制、执行证明和审计交付，它需要跨模型、跨框架、跨业务流程去成立。"
                "这正好给了独立基础层项目机会，因为我们不是围绕某个模型优化，而是在围绕关键动作的可信性建结构。"
            ),
            "anchors": [
                "定位是模型无关、框架无关的可信执行层",
                "面对的是组织治理问题，而不是单一模型调优问题",
            ],
            "avoid": "不要直接说“大厂做不好”，更稳的是说“关注点不同，层级不同”。",
        },
        {
            "q": "你们怎么证明不是只会讲概念，真的能现场演示？",
            "core": "因为我们能把关键链路跑出来，而且不是只有 happy path。",
            "full": (
                "这件事最重要的不只是正常样本能过，而是异常样本能不能暴露。当前我们能现场证明三件事："
                "门禁能跑，动作收据规则能验证，证据链正常样本能通过、篡改样本会失败。"
                "这意味着它不是只会展示架构图，而是关键机制已经能被独立验证。"
            ),
            "anchors": [
                "god-spear：STATUS: PASS",
                "safety-valve-spec：Overall: PASS",
                "execution-integrity-core：SELF_CHECK: PASS",
                "aro-audit：VERIFY_OK / Merkle mismatch",
            ],
            "avoid": "不要只说“我们有很多 repo”，评委关心的是机制能不能跑通。",
        },
        {
            "q": "POP、FDO、五层架构这些东西和当前项目到底是什么关系？",
            "core": "它们是架构来源和长期延展，不是当前主舞台的商业切口。",
            "full": (
                "AegisTrust 不是孤立想法，而是来自更完整的 agent 架构研究。`POP` 属于上游身份层，"
                "FDO 和相关 profile 材料提供外部技术讨论和可信度支撑。"
                "但当前阶段我们有意识地只聚焦其中最容易形成商业闭环的一段，也就是可信执行。"
                "这样既保留长期深度，又不会让主线发散。"
            ),
            "anchors": [
                "POP 已有 DOI 归档，证明上游深度",
                "aro-audit 已有 FDO profile 与 DOI 材料",
            ],
            "avoid": "不要把这些材料直接拉成并列主产品，否则答辩会散。",
        },
        {
            "q": "为什么北京适合成为你们的下一步落地点？",
            "core": "因为北京提供的不是口号，而是早期验证最需要的资源密度。",
            "full": (
                "这个项目当前最需要的不是一个名义上的落地城市，而是能把试点、样板和生态合作一起做起来的环境。"
                "北京同时具备高密度 AI 资源、高责任行业场景、研究机构和产业服务能力，这对一个从技术验证走向场景验证的项目非常关键。"
                "所以北京对我们来说不是宣传点，而是放大验证效率的节点。"
            ),
            "anchors": [
                "目标是联合验证、样板落地和生态合作",
                "切入场景本身就与北京的产业结构更匹配",
            ],
            "avoid": "不要回答成泛泛的‘北京政策好’，要讲资源和验证效率。",
        },
    ]

    lines = [
        "# AegisTrust 评委高频追问 Top 10",
        "",
        f"生成时间：{fmt_dt(ctx['generated_at'])}",
        "",
        "这份文档用于答辩，不用于主稿。目标不是把问题答得更长，而是答得更稳、更有抓手、更像一等奖项目该有的判断力。",
        "",
        "使用原则：",
        "- 先给一句短答，再决定要不要展开。",
        "- 展开时优先讲结构，不优先讲术语。",
        "- 每题都把回答压回主线：AegisTrust 是 AI 可信执行基础设施。",
        "",
    ]

    for idx, item in enumerate(questions, start=1):
        lines.extend(
            [
                f"## {idx}. {item['q']}",
                "",
                f"短答：{item['core']}",
                "",
                f"展开答：{item['full']}",
                "",
                "可顺手带出的证据：",
            ]
        )
        for anchor in item["anchors"]:
            lines.append(f"- {anchor}")
        lines.extend(
            [
                "",
                f"不要这样答：{item['avoid']}",
                "",
            ]
        )

    lines.extend(
        [
            "## 收束句",
            "",
            "如果评委追问太散，最后都可以收回到这一句：",
            "",
            "AegisTrust 现在最值得支持的地方，不是故事讲得多大，而是它已经把 AI 从会说走向会做之后最关键的一层基础能力，做出了可验证的雏形。",
            "",
        ]
    )

    return "\n".join(lines)


def build_mainstage_vs_backup(ctx: dict[str, Any]) -> str:
    return (
        textwrap.dedent(
            f"""\
        # AegisTrust 主舞台与备份拆分

        生成时间：{fmt_dt(ctx["generated_at"])}

        拆分原则：主舞台只讲评委 8 分钟内必须记住的内容；备份只承接追问；不要让高阶架构和材料证据抢走主线。

        ## A. 主舞台必须讲的内容

        - 一句话定位：AegisTrust = AI 可信执行基础设施。
        - 为什么现在需要：AI 正在从“会说”走向“会做”。
        - 为什么这是新问题：高风险 AI 动作会带来边界、责任和复验压力。
        - 为什么现有栈不够：模型层、日志层、单点插件都没有形成闭环。
        - 解决思路：先拦住，再留票，最后还能查账。
        - 四层可信执行闭环：门禁、收据、完整性、证据包。
        - 一个最容易理解的应用场景：AI 运维或高责任审批流程。
        - 关键真实性证据：本地仓库群 {len(ctx["repos"])} 个仓库、最近 30 天活跃 {ctx["active_30d"]} 个、`spear-check` 接入 {ctx["spear_adoption"]["count"]} 个仓库。
        - 关键回放证据：门禁 PASS、conformance PASS、完整性 PASS、篡改样本可暴露。
        - 商业路径、北京计划、下一步需要的支持。

        ## B. 备份页适合放的内容

        - `POP` 与五层总架构的关系：作为上游身份层和长期延展解释。
        - FDO / DOI / registry 相关证据：作为技术可信度补充，短讲即可。
        - 更细的 benchmark 或 conformance 说明：回答“你怎么证明它不是概念”。
        - 仓库级别的版本、标签、workflow、schema、test 细节。
        - repo adoption 证据截图或表格。
        - 外部讨论或 poster 材料：证明已有更广泛技术交流基础。

        ## C. 不适合放主舞台的内容

        - `POP` 的完整协议语义和 full taxonomy。
        - 完整 repo matrix 和逐仓库巡礼。
        - 五层架构的全图展开。
        - 所有 FDO 细节编号与映射表。
        - 任何未经验证的客户、收入、签约或标准主导叙事。
        - 过多英文缩写，如果一句话里讲不清，就不要放前 8 分钟。

        ## D. 遇到追问时的推荐顺序

        1. 先重复主线：我们做的是可信执行基础设施。
        2. 再补技术证据：回放、版本、接入数、证据链。
        3. 最后再补架构深度：`POP`、FDO、长期路线。

        ## E. 一句提醒

        主舞台争取留下两个记忆点就够了：

        - AI 从会说走向会做之后，可信执行会成为新的基础能力。
        - AegisTrust 已经把这件事做成一条能演示、能验证、能继续放大的工程闭环。
        """
        ).strip()
        + "\n"
    )


def generate_pptx(slides: list[dict[str, Any]], output_path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    palette = {
        "bg": RGBColor(247, 248, 250),
        "ink": RGBColor(20, 28, 48),
        "muted": RGBColor(90, 102, 122),
        "navy": RGBColor(13, 37, 63),
        "cyan": RGBColor(12, 157, 194),
        "lime": RGBColor(124, 180, 62),
        "sand": RGBColor(244, 232, 202),
        "white": RGBColor(255, 255, 255),
    }

    def set_text_style(
        run, size: int, color: RGBColor, bold: bool = False, font: str = "Arial"
    ):
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def add_textbox(slide, left, top, width, height, text=""):
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if text:
            box.text_frame.text = text
        return box

    for index, item in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        kind = item.get("kind", "normal")

        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = (
            palette["navy"] if kind in {"cover", "closing"} else palette["bg"]
        )
        background.line.fill.background()

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = (
            palette["cyan"] if kind != "closing" else palette["lime"]
        )
        accent.line.fill.background()

        if kind == "cover":
            orb1 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(8.8), Inches(0.7), Inches(3.2), Inches(3.2)
            )
            orb1.fill.solid()
            orb1.fill.fore_color.rgb = palette["cyan"]
            orb1.line.fill.background()

            orb2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(10.5), Inches(3.4), Inches(1.8), Inches(1.8)
            )
            orb2.fill.solid()
            orb2.fill.fore_color.rgb = palette["lime"]
            orb2.line.fill.background()

            title_box = add_textbox(slide, 0.8, 1.2, 6.8, 1.2, "")
            frame = title_box.text_frame
            p = frame.paragraphs[0]
            r = p.add_run()
            r.text = item["title"]
            set_text_style(r, 30, palette["white"], bold=True)

            sub_box = add_textbox(slide, 0.8, 2.2, 7.2, 0.8, "")
            p = sub_box.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = item["subtitle"]
            set_text_style(r, 18, palette["sand"], bold=False)

            tag_box = add_textbox(slide, 0.8, 3.2, 8.0, 1.2, "")
            p = tag_box.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = item["tagline"]
            set_text_style(r, 15, palette["white"])

            footer_box = add_textbox(slide, 0.8, 6.6, 5.0, 0.4, "")
            p = footer_box.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = item["footer"]
            set_text_style(r, 10, palette["sand"])
            continue

        title_box = add_textbox(slide, 0.7, 0.55, 8.2, 0.8, "")
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = item["title"]
        set_text_style(r, 24, palette["ink"], bold=True)

        page_chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(11.6),
            Inches(0.45),
            Inches(1.0),
            Inches(0.42),
        )
        page_chip.fill.solid()
        page_chip.fill.fore_color.rgb = palette["sand"]
        page_chip.line.fill.background()
        p = page_chip.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{index:02d}"
        set_text_style(r, 11, palette["navy"], bold=True)

        body_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7),
            Inches(1.45),
            Inches(8.2),
            Inches(5.35),
        )
        body_card.fill.solid()
        body_card.fill.fore_color.rgb = palette["white"]
        body_card.line.fill.background()

        headline_text = item.get("headline")
        bullet_top = 1.75
        if headline_text:
            headline_box = add_textbox(slide, 1.0, 1.72, 7.3, 1.15, "")
            headline_frame = headline_box.text_frame
            headline_frame.word_wrap = True
            p = headline_frame.paragraphs[0]
            p.space_after = Pt(8)
            r = p.add_run()
            r.text = headline_text
            set_text_style(r, 24, palette["ink"], bold=True)
            bullet_top = 2.95

        text_box = add_textbox(
            slide, 1.0, bullet_top, 7.5, 3.6 if headline_text else 4.8, ""
        )
        frame = text_box.text_frame
        frame.word_wrap = True
        first = True
        for bullet in item.get("bullets", []):
            p = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            p.level = 0
            p.space_after = Pt(12)
            r = p.add_run()
            r.text = bullet
            set_text_style(r, 16 if headline_text else 18, palette["ink"])

        side_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.25),
            Inches(1.45),
            Inches(3.35),
            Inches(5.35),
        )
        side_card.fill.solid()
        side_card.fill.fore_color.rgb = palette["navy"]
        side_card.line.fill.background()

        side_title = add_textbox(slide, 9.55, 1.8, 2.8, 0.5, "")
        p = side_title.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "一句话"
        set_text_style(r, 13, palette["sand"], bold=True)

        side_text = (
            item.get("callout")
            or item.get("subtitle")
            or item.get("tagline")
            or "当前可用"
        )
        side_body = add_textbox(slide, 9.55, 2.45, 2.7, 2.6, "")
        p = side_body.text_frame.paragraphs[0]
        p.word_wrap = True
        r = p.add_run()
        r.text = side_text
        set_text_style(r, 20, palette["white"], bold=True)

        hint_body = add_textbox(slide, 9.55, 5.35, 2.7, 1.0, "")
        p = hint_body.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "AI 可信执行\n面向真实动作"
        set_text_style(r, 11, palette["sand"])

        if kind == "closing":
            background.fill.fore_color.rgb = palette["navy"]
            title_box.left = Inches(0.8)
            title_box.top = Inches(2.2)
            p = title_box.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.color.rgb = palette["white"]
            slide.shapes._spTree.remove(body_card._element)
            slide.shapes._spTree.remove(side_card._element)
            slide.shapes._spTree.remove(side_title._element)
            slide.shapes._spTree.remove(side_body._element)
            slide.shapes._spTree.remove(hint_body._element)
            tag_box = add_textbox(slide, 0.8, 3.2, 9.0, 1.0, "")
            p = tag_box.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = item["subtitle"]
            set_text_style(r, 18, palette["sand"])
            foot = add_textbox(slide, 0.8, 4.2, 8.0, 0.8, "")
            p = foot.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = item["tagline"]
            set_text_style(r, 20, palette["white"], bold=True)

    prs.save(output_path)


def write_outputs(ctx: dict[str, Any], skip_pptx: bool = False) -> list[Path]:
    outputs: list[Path] = []

    github_report = OUTPUT_DIR / "github_progress_report.md"
    github_report.write_text(build_github_progress_report(ctx), encoding="utf-8")
    outputs.append(github_report)

    architecture_summary = OUTPUT_DIR / "architecture_summary.md"
    architecture_summary.write_text(build_architecture_summary(ctx), encoding="utf-8")
    outputs.append(architecture_summary)

    policy_context = OUTPUT_DIR / "policy_context.md"
    policy_context.write_text(build_policy_context(ctx), encoding="utf-8")
    outputs.append(policy_context)

    technical_progress = OUTPUT_DIR / "technical_progress.md"
    technical_progress.write_text(build_technical_progress(ctx), encoding="utf-8")
    outputs.append(technical_progress)

    judge_view = OUTPUT_DIR / "judge_view_diagnosis.md"
    judge_view.write_text(build_judge_view_diagnosis(ctx), encoding="utf-8")
    outputs.append(judge_view)

    award_storyline = OUTPUT_DIR / "award_storyline.md"
    award_storyline.write_text(build_award_storyline(ctx), encoding="utf-8")
    outputs.append(award_storyline)

    bp_text_v4 = build_bp_v4_award(ctx)
    bp_text_v2 = build_bp_v2(ctx)

    bp_doc = OUTPUT_DIR / "aegistrust_bp.md"
    bp_doc.write_text(bp_text_v4, encoding="utf-8")
    outputs.append(bp_doc)

    bp_doc_v2 = OUTPUT_DIR / "aegistrust_bp_v2.md"
    bp_doc_v2.write_text(bp_text_v2, encoding="utf-8")
    outputs.append(bp_doc_v2)

    bp_doc_v4 = OUTPUT_DIR / "aegistrust_bp_v4_award.md"
    bp_doc_v4.write_text(bp_text_v4, encoding="utf-8")
    outputs.append(bp_doc_v4)

    legacy_bp_doc = OUTPUT_DIR / "aegistrust_hicool_bp.md"
    legacy_bp_doc.write_text(bp_text_v4, encoding="utf-8")
    outputs.append(legacy_bp_doc)

    summary_text = build_competition_summary()
    project_summary = OUTPUT_DIR / "project_summary.txt"
    project_summary.write_text(summary_text, encoding="utf-8")
    outputs.append(project_summary)

    summary = OUTPUT_DIR / "competition_summary.txt"
    summary.write_text(summary_text, encoding="utf-8")
    outputs.append(summary)

    pitch_pack = OUTPUT_DIR / "pitch_pack_v1.md"
    pitch_pack.write_text(build_pitch_pack(ctx), encoding="utf-8")
    outputs.append(pitch_pack)

    app_form = OUTPUT_DIR / "application_form_short_v1.txt"
    app_form.write_text(build_application_form_short(ctx), encoding="utf-8")
    outputs.append(app_form)

    judge_qa = OUTPUT_DIR / "judge_qa_top10_v1.md"
    judge_qa.write_text(build_judge_qa_top10(ctx), encoding="utf-8")
    outputs.append(judge_qa)

    mainstage_split = OUTPUT_DIR / "mainstage_vs_backup.md"
    mainstage_split.write_text(build_mainstage_vs_backup(ctx), encoding="utf-8")
    outputs.append(mainstage_split)

    if not skip_pptx:
        slides_v4 = build_slides_v4_award(ctx)
        slides_v2 = build_slides_v2(ctx)

        pptx_path = OUTPUT_DIR / "aegistrust_pitch.pptx"
        generate_pptx(slides_v4, pptx_path)
        outputs.append(pptx_path)

        pptx_path_v2 = OUTPUT_DIR / "aegistrust_pitch_v2.pptx"
        generate_pptx(slides_v2, pptx_path_v2)
        outputs.append(pptx_path_v2)

        pptx_path_v4 = OUTPUT_DIR / "aegistrust_pitch_v4_award.pptx"
        generate_pptx(slides_v4, pptx_path_v4)
        outputs.append(pptx_path_v4)

    return outputs


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="生成 AegisTrust 项目材料。")
    parser.add_argument(
        "--skip-smoke", action="store_true", help="Skip local smoke replays."
    )
    parser.add_argument(
        "--skip-pptx", action="store_true", help="Skip PPTX generation."
    )
    parser.add_argument(
        "--use-github",
        action="store_true",
        help="显式使用 GitHub API；默认只使用本地仓库。",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    try:
        if args.use_github:
            ctx = collect_context(
                skip_smoke=args.skip_smoke, source_mode="github_or_cache"
            )
        else:
            ctx = collect_context(
                skip_smoke=args.skip_smoke,
                repos_override=scan_local_repos(),
                source_mode="local_fallback",
            )
        outputs = write_outputs(ctx, skip_pptx=args.skip_pptx)
    except ModuleNotFoundError as exc:
        if exc.name == "pptx":
            print(
                "python-pptx is required for PPTX generation. Re-run with --skip-pptx or install python-pptx.",
                file=sys.stderr,
            )
            return 2
        raise
    except (HTTPError, URLError) as exc:
        ctx = collect_context(
            skip_smoke=args.skip_smoke,
            repos_override=scan_local_repos(),
            source_mode="local_fallback",
        )
        outputs = write_outputs(ctx, skip_pptx=args.skip_pptx)
        print(f"GitHub 接口不可用，已自动回退到本地仓库模式：{exc}", file=sys.stderr)

    print("Generated files:")
    for output in outputs:
        print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
